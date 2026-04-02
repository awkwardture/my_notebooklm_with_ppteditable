#!/usr/bin/env python3
"""
从 PPTX 文件生成带缩略图的页面级模板。
每页导出为图片，并提取详细的风格描述。
使用 LibreOffice + pdftoppm 生成真实幻灯片截图。
"""

import os
import json
import glob
import subprocess
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

# PPT 样例目录
ppt_dir = None
# 尝试多个可能的路径
for search_path in [
    os.path.join(os.path.dirname(__file__), "ppt*"),
    os.path.join(os.path.dirname(os.path.dirname(__file__)), "ppt*"),
]:
    for item in glob.glob(search_path):
        if os.path.isdir(item) and not item.endswith('.py'):
            ppt_dir = item.strip()
            break
    if ppt_dir:
        break

if not ppt_dir:
    raise Exception("PPT 样例 directory not found")

output_dir = "/Users/wujianjun/data/VSC-PRJS/my_notebooklm_with_ppteditable/page_template"
thumbnail_dir = os.path.join(output_dir, "thumbnails")
os.makedirs(thumbnail_dir, exist_ok=True)

# 布局类型映射
LAYOUT_CATEGORY_MAP = {
    'title': '封面标题页',
    'content': '内容页',
    'table': '表格页',
    'chart': '图表页',
    'bullets': '列表页',
}


def convert_pptx_to_pdf(pptx_path, output_dir):
    """使用 LibreOffice 将 PPTX 转换为 PDF"""
    pptx_path = os.path.abspath(pptx_path).strip()
    output_dir = output_dir.strip()

    cmd = [
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        '--headless', '--convert-to', 'pdf',
        pptx_path, '--outdir', output_dir
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

    pdf_path = os.path.join(output_dir, os.path.basename(pptx_path).replace('.pptx', '.pdf'))
    return pdf_path if os.path.exists(pdf_path) else None


def extract_pdf_pages_to_images(pdf_path, output_prefix, size=(320, 180)):
    """使用 pdftoppm 将 PDF 每页转为图片"""
    pdf_path = pdf_path.strip()
    output_prefix = output_prefix.strip()

    cmd = ['pdftoppm', '-png', '-scale-to', str(size[0]), pdf_path, output_prefix]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

    # 返回生成的图片列表
    images = []
    idx = 1
    while True:
        img_path = f"{output_prefix}-{idx:02d}.png"
        if os.path.exists(img_path):
            images.append(img_path)
            idx += 1
        else:
            break
    return images


def analyze_slide(slide, slide_num, filename):
    """分析单页幻灯片"""
    elements = {
        'has_title': False,
        'has_subtitle': False,
        'has_text_boxes': 0,
        'has_bullets': False,
        'has_table': False,
        'table_structure': None,
        'has_chart': False,
        'chart_type': None,
        'has_image': False,
        'has_shape': False,
        'shape_count': 0,
    }

    # 分析标题
    title = slide.shapes.title if slide.shapes.title else None
    if title and title.text.strip():
        elements['has_title'] = True

    # 分析 shapes
    for shape in slide.shapes:
        try:
            if shape.has_table:
                elements['has_table'] = True
                elements['table_structure'] = {
                    'rows': len(shape.table.rows),
                    'cols': len(shape.table.columns),
                }
            elif shape.has_chart:
                elements['has_chart'] = True
                elements['chart_type'] = str(shape.chart.chart_type).split('.')[-1]
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                elements['has_image'] = True
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    elements['has_text_boxes'] += 1
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level > 0:
                            elements['has_bullets'] = True
                            break
            elif shape.shape_type in [1, 2, 3, 4, 5]:
                elements['has_shape'] = True
                elements['shape_count'] += 1
        except:
            continue

    # 确定布局类型
    if elements['has_title'] and not elements['has_table'] and not elements['has_chart'] and elements['has_text_boxes'] <= 2:
        layout_category = 'title'
    elif elements['has_table']:
        layout_category = 'table'
    elif elements['has_chart']:
        layout_category = 'chart'
    elif elements['has_bullets']:
        layout_category = 'bullets'
    else:
        layout_category = 'content'

    return elements, layout_category


def generate_style_description(filename, slide_num, elements, layout_category):
    """生成详细的风格描述 - 按天翼云企业级技术汇报标准格式"""

    # 基础风格模板
    style_template = """# PPT 样式风格描述

## 整体风格

- **定位**：企业级技术汇报 / 商务提案风格
- **基调**：专业、清晰、技术感、商务正式
- **气质**：兼具天翼云品牌调性与科技行业属性，传递可信赖的数字化服务形象

---

## 主色调

| 类型 | 色值 | 用途 |
|------|------|------|
| **主色** | #0066CC（天翼云蓝） | 标题、强调元素、主视觉 |
| **深蓝系** | #003366 / #336699 | 背景装饰、深色区块、网络与安全 |
| **浅蓝系** | #66CCFF / #0099FF | 基础设施、API 接口、背景点缀 |
| **辅助色** | #FFFFFF（白）、#666666（深灰） | 文字、对比元素 |
| **强调色** | #FF9933 / #FF6600（橙） | 痛点、挑战、流程环节 |
| **协作色** | #00AA55 / #00CC66（绿） | 积极信号、行动计划、协作沟通 |
| **科技色** | #9933FF（紫）、#00CCCC（青） | 平台组件、集成对接 |
| **商务色** | #FFAA00（金） | 预算、商务要素 |

---

## 字体风格

- **标题字体**：思源黑体（Source Han Sans）/ 思必达 / 方正兰亭黑
  - 粗体，字号 28-36pt，#0066CC 或深灰色
- **副标题字体**：中等字重，字号 18-24pt，#666666
- **正文字体**：思源黑体 Regular / 微软雅黑
  - 字号 14-18pt，行距 1.5 倍
- **数字/标签**：Roboto / DIN / 等宽字体，字号可适当缩小

---

## 背景风格

- **主背景**：纯白色 #FFFFFF 或浅灰 #F5F7FA
- **装饰背景**：
  - 抽象数字连接线条（细线，#66CCFF 半透明）
  - 云图标轮廓（简化几何形状，#E6F2FF 浅蓝）
  - 渐变蒙层（从左至右白→浅蓝，淡出效果）
- **分区背景**：深色区块（#003366）用于首屏标题区

---

## 图形元素

| 元素类型 | 风格描述 |
|----------|----------|
| **线条** | 细线风格（1-2pt），浅蓝或灰色，可带圆角端点 |
| **形状** | 扁平化几何形状：圆角矩形（卡片）、圆形（标签）、六边形（流程节点） |
| **图标** | 线性图标风格，2pt 线宽，圆角，与主题蓝同色系 |
| **图表** | 扁平化配色，数据区域填充半透明色，带圆角 |
| **连接线** | 虚线或点线表示间接关系，实线表示直接关联 |

---

## 排版风格

- **页面布局**：单栏或双栏，留白充足（上下边距 ≥ 60pt，左右边距 ≥ 40pt）
- **标题区**：顶部居中或左对齐，副标题紧随其后，留出呼吸空间
- **内容区**：卡片式布局或列表式布局，内容块间距 24-32pt
- **对齐方式**：左对齐为主，数字/标签可右对齐
- **留白比例**：40% 留白区域，确保信息层次清晰

---

## 本页视觉元素

**页面类型**：{layout_category_cn}

**来源**：{source_name}

**视觉构成**：
{visual_elements}

---

## 本页图表视觉建议

{chart_recommendations}

---

## 本页内容结构

{content_structure}

---

## 总结

本页应采用**{page_style_summary}**的视觉呈现，以天翼云蓝为主色调，通过{key_visual_method}可视化信息，确保内容层次分明、阅读流畅。
"""

    # 布局类型中文映射
    layout_map = {
        'title': '封面标题页',
        'content': '内容页',
        'table': '表格页',
        'chart': '图表页',
        'bullets': '列表页',
    }

    # 构建视觉元素描述
    visual_elements = []
    if elements['has_title']:
        visual_elements.append("- 顶部标题区（#0066CC 天翼云蓝，思源黑体 Bold 28-36pt）")
    if elements['has_table'] and elements['table_structure']:
        ts = elements['table_structure']
        visual_elements.append(f"- {ts['rows']}行×{ts['cols']}列表格（扁平化配色，半透明填充，圆角边框）")
    if elements['has_chart'] and elements['chart_type']:
        visual_elements.append(f"- {elements['chart_type']}图表（扁平化配色，数据标签清晰）")
    if elements['has_image']:
        visual_elements.append("- 装饰图片（科技风格，浅蓝云图标轮廓 #E6F2FF）")
    if elements['has_bullets']:
        visual_elements.append("- 项目符号列表（圆角图标，#0066CC 主题蓝，行距 1.5 倍）")
    if elements['has_text_boxes'] > 0:
        visual_elements.append(f"- {elements['has_text_boxes']} 个文本区域（卡片式布局，间距 24-32pt）")
    if elements['has_shape'] and elements['shape_count'] > 0:
        visual_elements.append(f"- {elements['shape_count']} 个装饰形状（扁平化几何图形，圆角矩形/圆形/六边形）")

    visual_elements_str = "\n".join(visual_elements) if visual_elements else "- 简洁内容布局"

    # 构建图表推荐
    chart_recommendations_map = {
        'title': "无需图表，重点突出标题视觉层次",
        'content': "可使用卡片式布局或列表式布局展示内容要点",
        'table': "表格采用扁平化配色，表头使用#0066CC，交替行背景色（#F5F7FA / #FFFFFF）增强可读性",
        'chart': "图表使用天翼云蓝系配色，主色#0066CC，辅助色#66CCFF，强调色#FF9933",
        'bullets': "列表项使用圆形图标（#0066CC），每项独立段落，行距 1.5 倍",
    }
    chart_recommendations = chart_recommendations_map.get(layout_category, "清晰展示内容要点")

    # 构建内容结构
    content_parts = []
    if elements['has_title']:
        content_parts.append("标题区（顶部，#0066CC）")
    if elements['has_table'] and elements['table_structure']:
        content_parts.append(f"表格区（{elements['table_structure']['rows']}×{elements['table_structure']['cols']}）")
    if elements['has_chart'] and elements['chart_type']:
        content_parts.append(f"图表区（{elements['chart_type']}）")
    if elements['has_text_boxes'] > 0:
        content_parts.append(f"内容区（{elements['has_text_boxes']}个文本块）")
    if elements['has_image']:
        content_parts.append("装饰图区")
    content_structure_str = " → ".join(content_parts) if content_parts else "单内容区布局"

    # 页面风格总结
    page_style_map = {
        'title': "简洁大气、视觉聚焦",
        'content': "清晰层次、信息密度适中",
        'table': "专业规整、数据清晰",
        'chart': "直观可视、重点突出",
        'bullets': "条理分明、易于扫读",
    }
    page_style_summary = page_style_map.get(layout_category, "专业清晰")

    # 关键视觉方法
    key_visual_map = {
        'title': "标题层级和视觉留白",
        'content': "卡片分组和图标引导",
        'table': "表格配色和行间距",
        'chart': "图表配色和数据标签",
        'bullets': "图标符号和段落间距",
    }
    key_visual_method = key_visual_map.get(layout_category, "内容结构化")

    return style_template.format(
        layout_category_cn=layout_map.get(layout_category, layout_category),
        source_name=filename.replace('.pptx', ''),
        visual_elements=visual_elements_str,
        chart_recommendations=chart_recommendations,
        content_structure=content_structure_str,
        page_style_summary=page_style_summary,
        key_visual_method=key_visual_method,
    )


# 收集所有页面模板
all_page_templates = []

# 获取所有 PPTX 文件（排除临时文件）
ppt_files = [f.strip() for f in os.listdir(ppt_dir) if f.endswith('.pptx') and not f.startswith('~')]
print(f"找到 {len(ppt_files)} 个 PPTX 文件")

# 批量转换 PPTX 为 PDF（并行处理）
pdf_cache = {}
temp_base = tempfile.mkdtemp(prefix="pptx_batch_")
print(f"临时目录：{temp_base}")

import concurrent.futures

def convert_single_pptx(filename):
    """转换单个 PPTX 为 PDF"""
    filepath = os.path.join(ppt_dir, filename).strip()
    ppt_temp_dir = os.path.join(temp_base, filename.replace('.pptx', ''))
    os.makedirs(ppt_temp_dir, exist_ok=True)

    pdf_path = convert_pptx_to_pdf(filepath, ppt_temp_dir)
    return (filename, pdf_path)

# 并行转换所有 PPTX 文件
print(f"\n=== 并行转换 PPTX 为 PDF（{len(ppt_files)}个文件）===")
with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
    futures = {executor.submit(convert_single_pptx, fn): fn for fn in ppt_files}
    for future in concurrent.futures.as_completed(futures):
        filename, pdf_path = future.result()
        if pdf_path:
            pdf_cache[filename] = pdf_path
            print(f"  ✓ {filename}: PDF 已生成")
        else:
            print(f"  ✗ {filename}: PDF 生成失败")

# 将 PDF 页面转为图片（并行处理）
print(f"\n=== 并行生成缩略图 ===")

def process_pptx_thumbnails(filename, pdf_path):
    """处理单个 PPTX 的缩略图生成"""
    style_name = filename.replace('.pptx', '').replace('副本', '')[:50].strip()
    img_prefix = os.path.join(os.path.dirname(pdf_path), f"thumb_{style_name}")

    images = extract_pdf_pages_to_images(pdf_path, img_prefix)
    results = []

    for idx, img_path in enumerate(images):
        slide_num = idx + 1
        thumb_filename = f"{style_name}_p{slide_num}.jpg"
        thumb_path = os.path.join(thumbnail_dir, thumb_filename)

        # 转换 PNG 为 JPEG
        from PIL import Image
        img = Image.open(img_path)
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        img.save(thumb_path, quality=85)

        # 删除临时 PNG
        os.remove(img_path)
        results.append((slide_num, thumb_path))

    return (filename, len(results), results)

# 并行生成所有缩略图
with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
    futures = {executor.submit(process_pptx_thumbnails, fn, pdf): (fn, pdf) for fn, pdf in pdf_cache.items()}
    for future in concurrent.futures.as_completed(futures):
        filename, count, results = future.result()
        print(f"  ✓ {filename}: {count} 页缩略图")

# 分析 PPTX 并生成模板数据（并行处理）
print(f"\n=== 并行分析幻灯片内容 ===")

def analyze_pptx_slides(filename):
    """分析单个 PPTX 的所有幻灯片"""
    filepath = os.path.join(ppt_dir, filename).strip()
    style_name = filename.replace('.pptx', '').replace('副本', '')[:50].strip()

    templates = []
    prs = Presentation(filepath)

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        # 分析幻灯片
        elements, layout_category = analyze_slide(slide, slide_num, filename)

        # 生成风格描述
        style_desc = generate_style_description(filename, slide_num, elements, layout_category)

        # 缩略图路径
        thumb_filename = f"{style_name}_p{slide_num}.jpg"
        thumb_path = os.path.join(thumbnail_dir, thumb_filename)
        thumb_relative = f"thumbnails/{thumb_filename}" if os.path.exists(thumb_path) else None

        # 构建页面模板
        page_template = {
            "id": f"{style_name}_page_{slide_num}",
            "source_file": filename,
            "source_name": style_name,
            "page_num": slide_num,
            "thumbnail": thumb_relative,
            "layout_category": layout_category,
            "layout_category_cn": LAYOUT_CATEGORY_MAP.get(layout_category, layout_category),
            "style_description": style_desc,
            "elements": elements,
            "colors": {
                "primary": "#0066CC",
                "secondary": "#66CCFF",
                "accent": "#FF9933",
            }
        }

        templates.append(page_template)

    return (filename, templates)

# 并行分析所有 PPTX 文件
all_page_templates = []
with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
    futures = {executor.submit(analyze_pptx_slides, fn): fn for fn in ppt_files}
    for future in concurrent.futures.as_completed(futures):
        filename, templates = future.result()
        all_page_templates.extend(templates)
        print(f"  ✓ {filename}: {len(templates)} 页模板已分析")

# 清理临时目录
import shutil
try:
    shutil.rmtree(temp_base)
except:
    pass

# 保存所有页面模板
output_file = os.path.join(output_dir, "page_templates.json")
with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(all_page_templates, f, ensure_ascii=False, indent=2)

print(f"\n========== 完成 ==========")
print(f"总计：{len(all_page_templates)} 页模板")
print(f"输出文件：{output_file}")
print(f"缩略图目录：{thumbnail_dir}")

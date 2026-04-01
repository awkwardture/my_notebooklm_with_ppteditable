"""Generate PPTX from infographic images using MiniMax to write python-pptx code."""

import os
import re
import traceback

from src.aliyun_client import generate_text_with_images as generate_text_aliyun
from src.prompts import PPT_CODE_GEN_SYSTEM_PROMPT, PPT_CODE_GEN_USER_PROMPT

DEFAULT_MODEL = "qwen3.5-plus"

# PPT 代码生成模型配置（阿里云百炼平台）
PPT_MODELS = {
    "qwen3.5-plus": {
        "name": "Qwen3.5-Plus",
    },
    "kimi-k2.5": {
        "name": "Kimi-K2.5",
    },
}

# ── Shared boilerplate embedded in every generated script ──
_SCRIPT_HEADER = '''\
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape
'''


def _extract_code(response_text: str) -> str:
    """Extract ONLY the build_slide function from the AI response."""
    # Step 1: extract from markdown code block if present
    match = re.search(r'```python\s*\n(.*?)```', response_text, re.DOTALL)
    code = match.group(1).strip() if match else response_text.strip()

    # Step 2: find the build_slide function and extract only it
    match = re.search(r'(def build_slide\(slide\):.*)', code, re.DOTALL)
    if not match:
        return code

    func_text = match.group(1)
    lines = func_text.split('\n')
    result = [lines[0]]  # "def build_slide(slide):"
    for line in lines[1:]:
        # Stop at the next top-level definition or non-indented code
        # (but allow blank lines and comments inside the function)
        if line and not line[0].isspace() and not line.startswith('#') and line.strip():
            break
        result.append(line)

    # Remove trailing blank lines
    while result and not result[-1].strip():
        result.pop()

    return '\n'.join(result)


def generate_slide_code(
    image_path: str,
    page_num: int,
    total_pages: int,
    model: str = DEFAULT_MODEL,
) -> str:
    """Send an infographic image to Vision model and get back python-pptx code."""
    user_prompt = PPT_CODE_GEN_USER_PROMPT.format(
        page_num=page_num,
        total_pages=total_pages,
    )

    # 所有 PPT 模型都使用阿里云 API
    response = generate_text_aliyun(
        model=model,
        system_prompt=PPT_CODE_GEN_SYSTEM_PROMPT,
        user_prompt=user_prompt,
        image_paths=[image_path],
    )

    return _extract_code(response)


def _make_pptx_script(build_func_codes: list[tuple[str, str]], output_path: str) -> str:
    """Assemble a full runnable script.

    build_func_codes: list of (func_name, func_body) pairs.
    """
    parts = [_SCRIPT_HEADER]
    parts.append(f'\nOUTPUT_PATH = r"{output_path}"\n')

    slide_calls = []
    for i, (func_name, code) in enumerate(build_func_codes):
        parts.append(f"\n# ── Slide {i + 1} ──\n")
        parts.append(code)
        parts.append("\n")
        slide_calls.append(
            f"s{i} = prs.slides.add_slide(prs.slide_layouts[6])\n"
            f"{func_name}(s{i})"
        )

    parts.append(f"""
# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

{chr(10).join(slide_calls)}
prs.save(OUTPUT_PATH)
""")
    return "\n".join(parts)


def _rename_func(code: str, new_name: str) -> str:
    """Rename build_slide -> new_name in the code."""
    return code.replace("def build_slide(slide)", f"def {new_name}(slide)")


def build_single_slide_pptx(slide_code: str, output_path: str) -> tuple[bool, str]:
    """Generate a single-slide PPTX from one slide's code. Returns (success, error)."""
    func_name = "build_slide_1"
    renamed = _rename_func(slide_code, func_name)
    script = _make_pptx_script([(func_name, renamed)], output_path)

    # Save assembled script for debugging (use _full suffix to avoid overwriting code file)
    script_path = output_path.replace(".pptx", "_full.py")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(script)

    return _exec_script(script)


def test_slide_code(slide_code: str) -> tuple[bool, str]:
    """Test if a slide code can be executed without errors.

    Returns (success, error_message).
    This is faster than building a full PPTX file.
    """
    # Apply patches first
    patched_code = _patch_common_errors(slide_code)

    # Use the same header as the full script (defines helper functions and colors)
    # but remove the 'from pptx.chart.data import CategoryChartData' line since we add it
    test_script = f'''
{_SCRIPT_HEADER}

# Create test presentation
prs = Presentation()
prs.slide_width = int(914400 * 13.333)
prs.slide_height = int(914400 * 7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Define the build function
{patched_code}

# Execute the build function
build_slide(slide)
'''
    try:
        exec(compile(test_script, "<test_slide>", "exec"), {"__builtins__": __builtins__})
        return True, ""
    except Exception:
        return False, traceback.format_exc()


def build_single_slide_pptx_with_retry(
    slide_code: str,
    output_path: str,
    max_retries: int = 3,
    regenerate_func: callable = None,
    regenerate_args: tuple = None,
) -> tuple[bool, str, str]:
    """Generate a single-slide PPTX with automatic retry on failure.

    Args:
        slide_code: The python-pptx code for the slide
        output_path: Where to save the PPTX file
        max_retries: Maximum number of retry attempts (default 3)
        regenerate_func: Optional function to regenerate code on failure
        regenerate_args: Arguments to pass to regenerate_func

    Returns:
        (success, error_message, final_code)
    """
    current_code = slide_code

    for attempt in range(max_retries):
        # First test the code
        test_ok, test_error = test_slide_code(current_code)

        if test_ok:
            # Code is valid, build the PPTX
            success, error = build_single_slide_pptx(current_code, output_path)
            if success:
                return True, "", current_code
            # If build failed despite test passing, return the error
            return False, error, current_code

        # Test failed, try to fix or regenerate
        print(f"[Slide test] Attempt {attempt + 1}/{max_retries} failed:")
        print(test_error[:500])  # Print first 500 chars of error

        if attempt < max_retries - 1 and regenerate_func is not None:
            # Try to regenerate the code
            print(f"[Slide test] Regenerating code...")
            try:
                current_code = regenerate_func(*regenerate_args)
            except Exception as e:
                print(f"[Slide test] Regeneration failed: {e}")
                return False, f"Regeneration failed: {e}", current_code
        else:
            # No more retries or no regenerate function
            return False, test_error, current_code

    return False, "Max retries exceeded", current_code


def build_full_pptx(slide_codes: dict[int, str], output_path: str) -> tuple[bool, str]:
    """Generate a full PPTX from multiple slide codes.

    slide_codes: {page_num: code_string} (1-based page numbers).
    Returns (success, error).
    """
    func_pairs = []
    for page_num in sorted(slide_codes.keys()):
        func_name = f"build_slide_{page_num}"
        renamed = _rename_func(slide_codes[page_num], func_name)
        func_pairs.append((func_name, renamed))

    script = _make_pptx_script(func_pairs, output_path)

    script_path = output_path.replace(".pptx", "_full.py")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(script)

    return _exec_script(script)


def _patch_common_errors(code: str) -> str:
    """Auto-fix common mistakes the AI makes in generated python-pptx code."""
    # Fix strings with embedded quotes - common in Chinese text
    # Pattern: .text = "content with "quotes" inside" causes syntax error
    lines = code.split('\n')
    fixed_lines = []
    for line in lines:
        # Check for .text = "..." or p.text = "..." patterns
        if re.search(r'[pf][_\d]*\.text\s*=\s*"', line):
            # Find positions of all ASCII quotes in the line
            quote_positions = [i for i, c in enumerate(line) if c == '"']
            if len(quote_positions) >= 3:  # 3+ quotes means embedded quotes
                # Get the text assignment part
                match = re.search(r'[pf][_\d]*\.text\s*=\s*', line)
                if match:
                    start_pos = match.end()
                    # Find first quote after .text =
                    first_quote = line.find('"', start_pos)
                    # Find last quote in line (the real closing delimiter)
                    last_quote = line.rfind('"')
                    if first_quote >= 0 and last_quote > first_quote:
                        # Content between first and last quote needs fixing
                        content = line[first_quote+1:last_quote]
                        # Escape all quotes in content (both ASCII and Chinese)
                        content = content.replace('\\', '\\\\')  # First escape existing backslashes
                        content = content.replace('"', '\\"')   # Escape ASCII quotes
                        content = content.replace('"', '\\"')   # Escape Chinese left quote
                        content = content.replace('"', '\\"')   # Escape Chinese right quote
                        line = line[:first_quote] + '"' + content + '"' + line[last_quote+1:]
        fixed_lines.append(line)
    code = '\n'.join(fixed_lines)

    # .line.background() -> .line.fill.background()
    code = re.sub(r'\.line\.background\(\)', '.line.fill.background()', code)
    # .line.no_fill() -> .line.fill.background()
    code = re.sub(r'\.line\.no_fill\(\)', '.line.fill.background()', code)
    # MSO_CONNECTOR not imported -> add import if used
    if 'MSO_CONNECTOR' in code and 'from pptx.enum.shapes import MSO_CONNECTOR' not in code:
        code = 'from pptx.enum.shapes import MSO_CONNECTOR\n' + code
    # Replace invalid MSO_SHAPE members with ROUNDED_RECTANGLE
    from pptx.enum.shapes import MSO_SHAPE
    _valid_shapes = set(MSO_SHAPE.__members__.keys())
    def _fix_shape(m):
        name = m.group(1)
        if name in _valid_shapes:
            return m.group(0)
        return f'MSO_SHAPE.ROUNDED_RECTANGLE'
    code = re.sub(r'MSO_SHAPE\.([A-Z_0-9]+)', _fix_shape, code)

    # Fix add_group_shape() - this API doesn't accept shape arguments
    code = re.sub(
        r'(\w+)\s*=\s*slide\.shapes\.add_group_shape\([^)]+\)',
        r'# REMOVED: group_shape not supported\n    pass',
        code
    )
    code = re.sub(
        r'slide\.shapes\.add_group_shape\([^)]+\)',
        r'pass  # group_shape not supported',
        code
    )

    # Fix axis_labels -> tick_labels
    code = re.sub(r'\.axis_labels\b', '.tick_labels', code)

    # Fix tick_labels.delete() - should use has_tick_labels = False on axis
    code = re.sub(r'\.tick_labels\.delete\(\)', '.has_tick_labels = False', code)

    # Fix line.fore_color -> line.color
    code = re.sub(r'\.line\.fore_color\b', '.line.color', code)

    # Add missing imports if used
    if 'MSO_ANCHOR' in code and 'from pptx.enum.text import MSO_ANCHOR' not in code:
        code = 'from pptx.enum.text import MSO_ANCHOR\n' + code
    if 'XL_LEGEND_POSITION' in code and 'from pptx.enum.chart import XL_LEGEND_POSITION' not in code:
        code = 'from pptx.enum.chart import XL_LEGEND_POSITION\n' + code

    # Fix MsoArrowheadLength/MsoArrowheadWidth - these don't exist in python-pptx
    # Comment out arrowhead settings
    code = re.sub(
        r'(\s*)\w+\.end_arrowhead\.[^\n]+',
        r'\1# arrowhead settings not supported',
        code
    )
    code = re.sub(
        r'(\s*)\w+\.start_arrowhead\.[^\n]+',
        r'\1# arrowhead settings not supported',
        code
    )

    # Fix p.add_run("text") - add_run() takes no arguments in python-pptx
    # Correct: run = p.add_run(); run.text = "content"
    def _fix_add_run(m):
        indent = m.group(1)
        var = m.group(2)
        text = m.group(3)
        return f'{indent}{var} = p.add_run()\n{indent}{var}.text = "{text}"'
    code = re.sub(
        r'(\s*)(\w+)\s*=\s*p[\d_]*\.add_run\("([^"]*)"\)',
        _fix_add_run,
        code
    )

    # Always replace MSO_DASH_STYLE with MSO_LINE_DASH_STYLE (correct name)
    code = code.replace('MSO_DASH_STYLE', 'MSO_LINE_DASH_STYLE')

    # Fix MSO_SHAPE_TYPE incorrect usage
    code = code.replace('MSO_SHAPE_TYPE.STRAIGHT', 'MSO_CONNECTOR.STRAIGHT')
    code = re.sub(r'MSO_SHAPE_TYPE\.[A-Z_]+', 'MSO_CONNECTOR.STRAIGHT', code)

    # Fix add_connector with invalid types (MSO_SHAPE, MSO_CONNECTOR_TYPE, etc.)
    # Valid connector types: MSO_CONNECTOR.STRAIGHT, ELBOW, ELBOW_ARROW
    # First, replace MSO_SHAPE.LINE specifically (common mistake)
    code = code.replace('MSO_SHAPE.LINE', 'MSO_CONNECTOR.STRAIGHT')
    code = re.sub(
        r'add_connector\(\s*MSO_SHAPE\.[A-Z_]+',
        'add_connector(MSO_CONNECTOR.STRAIGHT',
        code
    )
    code = re.sub(
        r'add_connector\(\s*MSO_CONNECTOR_TYPE\.[A-Z_]+',
        'add_connector(MSO_CONNECTOR.STRAIGHT',
        code
    )
    code = re.sub(
        r'add_connector\(\s*MSO_AUTO_SHAPE_TYPE\.[A-Z_]+',
        'add_connector(MSO_CONNECTOR.STRAIGHT',
        code
    )

    # Fix incorrect tuple unpacking in enumerate: for j, (a, b) in enumerate(row):
    # This pattern is wrong when row is a flat list, not a list of tuples
    # Fix: for j, item in enumerate(row): a, b = item
    def _fix_enumerate_unpack(m):
        indent = m.group(1)
        var1 = m.group(2)
        var2 = m.group(3)
        var3 = m.group(4)
        row_var = m.group(5)
        return f'{indent}for {var1}, item in enumerate({row_var}):{indent}    {var2}, {var3} = item'
    code = re.sub(
        r'(\s*)for\s+(\w+),\s*\((\w+),\s*(\w+)\)\s+in\s+enumerate\((\w+)\):',
        _fix_enumerate_unpack,
        code
    )

    # Fix p.bullet usage - python-pptx paragraphs don't have .bullet attribute
    # Remove/comment out any .bullet related code
    code = re.sub(
        r'(\s*)[pf](_\w+)?\.bullet\.[^\n]+',
        r'\1# bullet not supported in python-pptx',
        code
    )
    code = re.sub(
        r'(\s*)[pf](_\w+)?\.bullet\s*=[^\n]+',
        r'\1# bullet not supported in python-pptx',
        code
    )

    # Fix add_line - python-pptx doesn't have add_line, use add_connector or add_shape
    # Convert add_line(x1, y1, x2, y2) to add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2-x1, y2-y1)
    def _fix_add_line(m):
        indent = m.group(1)
        var = m.group(2) if m.group(2) else ''
        x1 = m.group(3)
        y1 = m.group(4)
        x2 = m.group(5)
        y2 = m.group(6)
        if var:
            return f"{indent}{var} = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, {x1}, {y1}, {x2}, {y2})"
        else:
            return f"{indent}slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, {x1}, {y1}, {x2}, {y2})"
    code = re.sub(
        r'(\s*)(\w+\s*=\s*)?slide\.shapes\.add_line\(([^,]+),\s*([^,]+),\s*([^,]+),\s*([^)]+)\)',
        _fix_add_line,
        code
    )

    # Fix shape.adjustments - commenting out since it often causes IndexError
    # Only certain shapes (like rounded rectangle) have adjustments
    # Need to add pass statement to avoid empty try blocks
    def _fix_adjustments(m):
        indent = m.group(1)
        return f'{indent}# adjustments not supported\n{indent}pass'
    code = re.sub(
        r'(\s*)\w+\.adjustments\[\d+\]\s*=[^\n]+',
        _fix_adjustments,
        code
    )

    # Fix .fill.fore_color -> .fill.fore_color (correct)
    # No change needed, already correct

    return code


def _exec_script(script: str) -> tuple[bool, str]:
    """Execute a generated pptx script in-process. Returns (success, error)."""
    script = _patch_common_errors(script)
    try:
        exec(compile(script, "<pptx_gen>", "exec"), {"__builtins__": __builtins__})
        return True, ""
    except Exception:
        return False, traceback.format_exc()


# ── Slide code persistence ──

def get_slides_dir(proj_dir: str) -> str:
    d = os.path.join(proj_dir, "最终文档", "ppt_slides")
    os.makedirs(d, exist_ok=True)
    return d


def save_slide_code(proj_dir: str, page_num: int, code: str):
    path = os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.py")
    with open(path, "w", encoding="utf-8") as f:
        f.write(code)


def load_slide_code(proj_dir: str, page_num: int) -> str | None:
    path = os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.py")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    return None


def load_all_slide_codes(proj_dir: str) -> dict[int, str]:
    """Return {page_num: code} for all saved slide codes."""
    slides_dir = get_slides_dir(proj_dir)
    result = {}
    for fname in sorted(os.listdir(slides_dir)):
        m = re.match(r"slide_(\d+)\.py$", fname)
        if m:
            page_num = int(m.group(1))
            with open(os.path.join(slides_dir, fname), "r", encoding="utf-8") as f:
                result[page_num] = f.read()
    return result


def get_single_pptx_path(proj_dir: str, page_num: int) -> str:
    return os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.pptx")

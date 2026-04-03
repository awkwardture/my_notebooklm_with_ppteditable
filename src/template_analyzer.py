"""PPT 模板分析器模块

功能：
1. 上传 PPT 文件，提取每一页的缩略图
2. 使用 AI Vision 模型分析每一页的视觉风格
3. 生成 style_description 和 layout_category
4. 保存为模板格式
"""

import os
import json
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from io import BytesIO

from src.aliyun_client import generate_text_with_images
from src.prompts import PPT_TEMPLATE_ANALYSIS_SYSTEM_PROMPT, PPT_TEMPLATE_ANALYSIS_USER_PROMPT


# 布局类型映射
LAYOUT_CATEGORY_MAP = {
    "title": "封面标题页",
    "content": "内容页",
    "table": "表格页",
    "chart": "图表页",
    "bullets": "列表页",
}

LAYOUT_CATEGORY_MAP_REVERSE = {v: k for k, v in LAYOUT_CATEGORY_MAP.items()}


def extract_slide_thumbnails(pptx_path: str, output_dir: str) -> List[Dict[str, Any]]:
    """从 PPTX 文件中提取每一页的缩略图

    Args:
        pptx_path: PPTX 文件路径
        output_dir: 输出目录

    Returns:
        缩略图信息列表，每项包含 page_num, thumbnail_path, width, height
    """
    prs = Presentation(pptx_path)
    slides_info = []

    os.makedirs(output_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        page_num = i + 1

        # 使用 slide 的 thumbnail 方法（如果可用）
        # 否则生成一个简单的占位图
        try:
            # 创建一个临时图片来表示这一页
            # 由于 python-pptx 不直接支持导出图片，我们需要用其他方式
            # 这里先保存为占位符，实际使用时需要外部工具转换
            thumbnail_path = os.path.join(output_dir, f"p{page_num}.jpg")

            # 使用 pptx 的 shapes 信息生成一个简单的描述
            # 实际的缩略图需要通过其他方式生成（如 LibreOffice 或 Aspose）
            # 这里先创建一个空的占位文件
            with open(thumbnail_path, 'wb') as f:
                # 创建一个 1x1 的空白图片作为占位
                # 实际应该生成真实的缩略图
                f.write(b'')  # 占位，实际实现需要图片生成库

            slides_info.append({
                "page_num": page_num,
                "thumbnail_path": thumbnail_path,
                "width": prs.slide_width,
                "height": prs.slide_height,
            })
        except Exception as e:
            print(f"Error extracting thumbnail for slide {page_num}: {e}")
            slides_info.append({
                "page_num": page_num,
                "thumbnail_path": None,
                "width": prs.slide_width,
                "height": prs.slide_height,
            })

    return slides_info


def analyze_slide_style(
    image_path: str,
    page_num: int,
    model: str = "qwen3.5-plus"
) -> Dict[str, Any]:
    """使用 AI Vision 模型分析单页幻灯片的风格

    Args:
        image_path: 幻灯片缩略图路径
        page_num: 页码
        model: 使用的 Vision 模型

    Returns:
        分析结果，包含 layout_category, style_description, elements, colors 等
    """
    user_prompt = PPT_TEMPLATE_ANALYSIS_USER_PROMPT.format(page_num=page_num)

    try:
        response = generate_text_with_images(
            model=model,
            system_prompt=PPT_TEMPLATE_ANALYSIS_SYSTEM_PROMPT,
            user_prompt=user_prompt,
            image_paths=[image_path]
        )

        # 解析 JSON 响应
        import re
        json_match = re.search(r'```(?:json)?\s*(.*?)\s*```', response, re.DOTALL)
        if json_match:
            json_str = json_match.group(1).strip()
        else:
            json_str = response.strip()

        result = json.loads(json_str)
        result["page_num"] = page_num
        return result

    except Exception as e:
        print(f"Error analyzing slide {page_num}: {e}")
        return {
            "page_num": page_num,
            "layout_category": "content",
            "layout_category_cn": "内容页",
            "style_description": f"AI 分析失败：{str(e)}",
            "elements": {},
            "colors": {}
        }


def analyze_ppt_template(
    pptx_path: str,
    thumbnail_dir: Optional[str] = None,
    output_json_path: Optional[str] = None,
    source_name: Optional[str] = None,
    model: str = "qwen3.5-plus"
) -> List[Dict[str, Any]]:
    """分析整个 PPT 文件，为每一页生成风格模板

    Args:
        pptx_path: PPTX 文件路径
        thumbnail_dir: 缩略图输出目录（默认在当前目录创建 thumbnails 子目录）
        output_json_path: 输出 JSON 文件路径（可选，用于保存结果）
        source_name: 源文件名称（用于模板标识）
        model: 使用的 Vision 模型

    Returns:
        模板数据列表
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    # 获取文件名作为 source_name
    if not source_name:
        source_name = os.path.splitext(os.path.basename(pptx_path))[0]

    # 设置缩略图目录
    if not thumbnail_dir:
        base_dir = os.path.dirname(pptx_path)
        thumbnail_dir = os.path.join(base_dir, "thumbnails")

    # 由于 python-pptx 不直接支持导出图片，我们需要一个替代方案
    # 当前实现要求用户先提供缩略图，或使用外部工具生成
    # 这里我们返回一个需要外部生成缩略图的提示

    print(f"分析 PPT 文件：{pptx_path}")
    print(f"缩略图目录：{thumbnail_dir}")

    # 检查是否有预先生成的缩略图
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)

    templates = []
    for i in range(total_slides):
        page_num = i + 1
        thumbnail_path = os.path.join(thumbnail_dir, f"p{page_num}.jpg")

        if os.path.exists(thumbnail_path):
            # 使用现有缩略图分析
            print(f"分析第 {page_num} 页...")
            slide_analysis = analyze_slide_style(thumbnail_path, page_num, model)
        else:
            # 没有缩略图，使用占位分析
            print(f"警告：缩略图不存在 {thumbnail_path}，使用占位分析")
            slide_analysis = {
                "page_num": page_num,
                "layout_category": "content",
                "layout_category_cn": "内容页",
                "style_description": "待分析 - 需要生成缩略图后重新分析",
                "elements": {},
                "colors": {}
            }

        # 构建完整的模板数据
        template_data = {
            "id": f"{source_name}_page_{page_num}_{uuid.uuid4().hex[:8]}",
            "source_file": os.path.basename(pptx_path),
            "source_name": source_name,
            "page_num": page_num,
            "thumbnail": f"thumbnails/{source_name}_p{page_num}.jpg",
            "layout_category": slide_analysis.get("layout_category", "content"),
            "layout_category_cn": slide_analysis.get("layout_category_cn", "内容页"),
            "style_description": slide_analysis.get("style_description", ""),
            "elements": slide_analysis.get("elements", {}),
            "colors": slide_analysis.get("colors", {})
        }

        templates.append(template_data)

    # 保存结果
    if output_json_path:
        os.makedirs(os.path.dirname(output_json_path), exist_ok=True)
        # 如果文件已存在，加载并追加
        if os.path.exists(output_json_path):
            with open(output_json_path, 'r', encoding='utf-8') as f:
                existing = json.load(f)
            existing.extend(templates)
            templates = existing
        with open(output_json_path, 'w', encoding='utf-8') as f:
            json.dump(templates, f, ensure_ascii=False, indent=2)
        print(f"模板数据已保存到：{output_json_path}")

    return templates


def add_template_to_library(
    template_data: List[Dict[str, Any]],
    library_path: str = None
) -> bool:
    """将分析后的模板添加到模板库

    Args:
        template_data: 模板数据列表
        library_path: 模板库文件路径（默认为 page_template/page_templates.json）

    Returns:
        是否成功
    """
    if not library_path:
        base_dir = os.path.dirname(os.path.dirname(__file__))
        library_path = os.path.join(base_dir, "page_template", "page_templates.json")

    # 加载现有模板库
    existing_templates = []
    if os.path.exists(library_path):
        with open(library_path, 'r', encoding='utf-8') as f:
            existing_templates = json.load(f)

    # 检查是否有重复（根据 id）
    existing_ids = {t.get("id") for t in existing_templates}
    new_templates = []
    for t in template_data:
        if t.get("id") not in existing_ids:
            new_templates.append(t)
        else:
            print(f"跳过重复模板：{t.get('id')}")

    if not new_templates:
        print("没有新模板需要添加")
        return False

    # 追加新模板
    existing_templates.extend(new_templates)

    # 保存
    os.makedirs(os.path.dirname(library_path), exist_ok=True)
    with open(library_path, 'w', encoding='utf-8') as f:
        json.dump(existing_templates, f, ensure_ascii=False, indent=2)

    print(f"成功添加 {len(new_templates)} 个模板到模板库")
    return True


def pptx_to_thumbnails(pptx_path: str, output_dir: str) -> List[str]:
    """将 PPTX 文件转换为缩略图序列

    使用 LibreOffice 或类似工具将 PPTX 转换为图片
    需要系统安装 LibreOffice

    Args:
        pptx_path: PPTX 文件路径
        output_dir: 输出目录

    Returns:
        生成的缩略图路径列表
    """
    import subprocess

    os.makedirs(output_dir, exist_ok=True)

    # 检查 LibreOffice 是否可用
    try:
        result = subprocess.run(
            ["soffice", "--version"],
            capture_output=True,
            text=True,
            timeout=5
        )
        has_libreoffice = result.returncode == 0
    except:
        has_libreoffice = False

    if has_libreoffice:
        # 使用 LibreOffice 转换
        # soffice --headless --convert-to jpg --outdir /path/to/output file.pptx
        try:
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to", "jpg:JPEG_Thumb",
                    "--outdir", output_dir,
                    pptx_path
                ],
                timeout=120
            )
            # 获取生成的文件
            thumbnails = sorted([
                os.path.join(output_dir, f)
                for f in os.listdir(output_dir)
                if f.endswith('.jpg')
            ])
            return thumbnails
        except Exception as e:
            print(f"LibreOffice 转换失败：{e}")

    # 如果没有 LibreOffice，返回空列表，让用户手动处理
    print("警告：系统未安装 LibreOffice，无法自动生成缩略图")
    print("请手动安装 LibreOffice 后重试，或手动提供缩略图")
    return []


if __name__ == "__main__":
    # 示例用法
    import sys

    if len(sys.argv) > 1:
        pptx_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None

        # 先生成缩略图
        thumb_dir = os.path.join(os.path.dirname(pptx_file), "thumbnails")
        thumbnails = pptx_to_thumbnails(pptx_file, thumb_dir)

        if thumbnails:
            print(f"生成了 {len(thumbnails)} 张缩略图")

            # 分析 PPT
            templates = analyze_ppt_template(
                pptx_path=pptx_file,
                thumbnail_dir=thumb_dir,
                output_json_path=output_file
            )

            print(f"分析了 {len(templates)} 页")
        else:
            print("无法生成缩略图，请确保系统安装了 LibreOffice")
    else:
        print("Usage: python template_analyzer.py <pptx_file> [output_json]")

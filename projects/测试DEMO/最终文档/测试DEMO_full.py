from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"projects/测试DEMO/最终文档/测试DEMO.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BLUE_DARK = RGBColor(0x1F, 0x4E, 0x79)
    BLUE_PRIMARY = RGBColor(0x00, 0x70, 0xC0)
    BLUE_LIGHT_BG = RGBColor(0xF4, 0xF8, 0xFC)
    BLUE_BORDER = RGBColor(0xBD, 0xD7, 0xEE)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    GRAY_BAR = RGBColor(0xA6, 0xA6, 0xA6)
    ORANGE_WARN = RGBColor(0xED, 0x7D, 0x31)
    BLACK_TEXT = RGBColor(0x33, 0x33, 0x33)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "简单场景：AI 初显 50% 提效潜力 ⚡"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_NAME

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "案例A —— 标准 ALV 销售订单报表验证"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = FONT_NAME

    # 3. Page Indicator (Top Right)
    page_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.8), Inches(0.6), Inches(1.0), Inches(0.4))
    page_shape.fill.solid()
    page_shape.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    page_shape.line.fill.background()
    tf = page_shape.text_frame
    tf.text = "P1 / 4"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = BLUE_PRIMARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 4. Left Column - Core Points Heading
    heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3), Inches(0.5))
    tf = heading_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心要点"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = FONT_NAME

    # 5. Left Column - Core Points Items
    items = [
        ("⚡", "1. 效率突破: ", "Claude Code 仅需 30 分钟完成开发，较手写提效 50%。"),
        ("🤖", "2. 工具差距: ", "GitHub Copilot 耗时与手写持平，且因 SQL 不兼容导致运行崩溃。"),
        ("🔀", "3. 交互对比: ", "Claude 仅需 2 轮提示即可运行，Copilot 需 5 轮以上人工干预。"),
        ("🔍", "4. 核心瓶颈: ", "[地址取数逻辑] 与 [过账状态字段] 是 AI 初始生成的共同盲区。")
    ]

    y_offset = 2.7
    for i, (icon, label, desc) in enumerate(items):
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(0.6), Inches(0.6))
        icon_box.text_frame.text = icon
        icon_box.text_frame.paragraphs[0].font.size = Pt(24)

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.2), y_offset, Inches(4.8), Inches(0.8))
        text_box.text_frame.word_wrap = True
        p = text_box.text_frame.paragraphs[0]
        
        run1 = p.add_run()
        run1.text = label
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = BLUE_DARK
        run1.font.name = FONT_NAME
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(14)
        run2.font.color.rgb = BLACK_TEXT
        run2.font.name = FONT_NAME

        # Separator line
        if i < len(items) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), y_offset + 0.85, Inches(6.0), y_offset + 0.85)
            line.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        
        y_offset += 1.0

    # 6. Right Column - Chart Box Background
    chart_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.3), Inches(4.2))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = BLUE_LIGHT_BG
    chart_bg.line.color.rgb = BLUE_BORDER
    chart_bg.line.width = Pt(1.5)

    # Chart Title
    chart_title = slide.shapes.add_textbox(Inches(6.7), Inches(2.0), Inches(4), Inches(0.5))
    tf = chart_title.text_frame
    p = tf.paragraphs[0]
    p.text = "开发耗时对比 (分钟)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = FONT_NAME

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ["手写/Copilot\n(Manual/AI)", "Claude Code\n(AI)"]
    chart_data.add_series('耗时', [60, 30])

    x, y, cx, cy = Inches(6.6), Inches(2.8), Inches(5.8), Inches(2.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.plots[0].has_data_labels = False # We will add custom text boxes for labels
    chart.plots[0].gap_width = 100

    # Customize bar colors
    series = chart.plots[0].series[0]
    pt0 = series.points[0]
    pt0.format.fill.solid()
    pt0.format.fill.fore_color.rgb = GRAY_BAR
    pt1 = series.points[1]
    pt1.format.fill.solid()
    pt1.format.fill.fore_color.rgb = BLUE_PRIMARY

    # Custom Data Labels
    lbl1 = slide.shapes.add_textbox(Inches(11.9), Inches(4.15), Inches(1), Inches(0.4))
    lbl1.text_frame.text = "⚠️ 60"
    lbl1.text_frame.paragraphs[0].font.size = Pt(14)
    lbl1.text_frame.paragraphs[0].font.name = FONT_NAME

    lbl2 = slide.shapes.add_textbox(Inches(9.6), Inches(3.25), Inches(1), Inches(0.4))
    lbl2.text_frame.text = "⚡ 30"
    lbl2.text_frame.paragraphs[0].font.size = Pt(14)
    lbl2.text_frame.paragraphs[0].font.name = FONT_NAME

    # Chart Callout (50% 提效)
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, Inches(9.1), Inches(2.4), Inches(1.6), Inches(0.6))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    callout.line.color.rgb = BLUE_BORDER
    tf = callout.text_frame
    p = tf.paragraphs[0]
    p.text = "50% 提效"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Chart Warning Text
    warn_text = slide.shapes.add_textbox(Inches(10.3), Inches(4.6), Inches(2.5), Inches(0.4))
    tf = warn_text.text_frame
    p = tf.paragraphs[0]
    p.text = "SQL 不兼容/运行崩溃"
    p.font.size = Pt(10)
    p.font.color.rgb = ORANGE_WARN
    p.font.name = FONT_NAME

    # 7. Right Column - Conclusion Box
    conc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(6.2), Inches(6.3), Inches(0.8))
    conc_bg.fill.solid()
    conc_bg.fill.fore_color.rgb = BLUE_LIGHT_BG
    conc_bg.line.color.rgb = BLUE_BORDER
    conc_bg.line.width = Pt(1.5)

    conc_text = slide.shapes.add_textbox(Inches(6.5), Inches(6.35), Inches(6.3), Inches(0.5))
    tf = conc_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run1 = p.add_run()
    run1.text = "✔️ ⚡ AI 初显 "
    run1.font.size = Pt(22)
    run1.font.bold = True
    run1.font.color.rgb = BLACK_TEXT
    run1.font.name = FONT_NAME

    run2 = p.add_run()
    run2.text = "50% 提效潜力"
    run2.font.size = Pt(22)
    run2.font.bold = True
    run2.font.color.rgb = BLUE_PRIMARY
    run2.font.name = FONT_NAME



# ── Slide 2 ──

def build_slide_2(slide):
    # 定义颜色
    BG_COLOR = RGBColor(43, 43, 43)
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(170, 170, 170)
    ORANGE = RGBColor(243, 139, 0)
    PANEL_BG = RGBColor(56, 56, 56)
    PANEL_BORDER = RGBColor(80, 80, 80)
    CHART_GRAY = RGBColor(120, 120, 120)
    DARK_ORANGE_BG = RGBColor(80, 60, 40)
    DARK_GRAY_BG = RGBColor(65, 65, 65)

    # 设置背景色
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR

    # 1. 标题区域
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "中等难度：陷入“幻觉”与语法泥潭"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.6))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "案例B — 采购配额维护（函数组开发）"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"

    # 2. 左侧：图表区域
    # 图表标题
    chart_title = slide.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(2.5), Inches(0.4))
    p = chart_title.text_frame.paragraphs[0]
    p.text = "错误分布饼图"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 环形图
    chart_data = CategoryChartData()
    chart_data.categories = ['虚构字段', '其他错误']
    chart_data.add_series('Series 1', (50, 50))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(2.0), Inches(2.3), Inches(2.5), Inches(2.5), chart_data).chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = False
    
    # 设置图表颜色
    points = chart.plots[0].series[0].points
    points[0].format.fill.solid()
    points[0].format.fill.fore_color.rgb = ORANGE
    points[1].format.fill.solid()
    points[1].format.fill.fore_color.rgb = CHART_GRAY

    # 图表中心图标 (带问号的警告三角)
    center_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.05), Inches(3.3), Inches(0.4), Inches(0.35))
    center_tri.fill.background()
    center_tri.line.color.rgb = ORANGE
    center_tri.line.width = Pt(1.5)
    
    center_text = slide.shapes.add_textbox(Inches(3.05), Inches(3.35), Inches(0.4), Inches(0.35))
    p = center_text.text_frame.paragraphs[0]
    p.text = "?"
    p.font.color.rgb = ORANGE
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    center_text.text_frame.margin_top = Pt(2)

    # 图表标签 - 右侧 (橙色)
    lbl_r = slide.shapes.add_textbox(Inches(4.6), Inches(2.5), Inches(2.0), Inches(1.0))
    tf = lbl_r.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "虚构字段\n(AI Hallucinations)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "50%"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p3 = tf.add_paragraph()
    p3.text = "9处核心字段"
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY_TEXT

    # 图表标签 - 左侧 (灰色)
    lbl_l = slide.shapes.add_textbox(Inches(0.2), Inches(3.1), Inches(1.8), Inches(1.0))
    tf = lbl_l.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "其他错误\n(Other Errors)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.RIGHT
    p2 = tf.add_paragraph()
    p2.text = "21个连锁语法错误"
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.RIGHT

    # 连接线
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.8), Inches(3.4), Inches(2.1), Inches(3.4))
    line1.line.color.rgb = GRAY_TEXT
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.4), Inches(3.4), Inches(4.7), Inches(3.4))
    line2.line.color.rgb = ORANGE

    # 3. 左侧：对比表格区域
    table_top = Inches(5.2)
    table_left = Inches(0.5)
    table_width = Inches(5.5)
    table_height = Inches(2.0)

    # 外框面板
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left, table_top, table_width, table_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG
    panel.line.color.rgb = ORANGE
    panel.line.width = Pt(1)

    # 面板标题
    p_title = slide.shapes.add_textbox(table_left, table_top + Inches(0.05), table_width, Inches(0.3))
    p = p_title.text_frame.paragraphs[0]
    p.text = "AI 虚构字段 vs SAP 实际字段"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 表头
    h_y = table_top + Inches(0.4)
    h_w = Inches(2.4)
    h_h = Inches(0.25)
    
    # 左表头
    h_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(0.2), h_y, h_w, h_h)
    h_left.fill.solid()
    h_left.fill.fore_color.rgb = DARK_ORANGE_BG
    h_left.line.color.rgb = ORANGE
    p = h_left.text_frame.paragraphs[0]
    p.text = "AI 虚构"
    p.font.color.rgb = ORANGE
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 右表头
    h_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(2.9), h_y, h_w, h_h)
    h_right.fill.solid()
    h_right.fill.fore_color.rgb = PANEL_BORDER
    h_right.line.color.rgb = CHART_GRAY
    p = h_right.text_frame.paragraphs[0]
    p.text = "SAP 实际"
    p.font.color.rgb = WHITE
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 数据行
    rows_data = [
        ("Z_QUOTA_QTY", "MENG"),
        ("Z_SOURCE_VEND", "LIFNR"),
        ("Z_QUOTA_UNIT", "MEINS"),
        ("Z_VALID_TO", "DATBI"),
        ("Z_QUOTA_TYPE", "QUNUM")
    ]
    row_y_start = h_y + Inches(0.3)
    row_h = Inches(0.22)
    row_spacing = Inches(0.26)

    for i, (ai_field, sap_field) in enumerate(rows_data):
        y = row_y_start + i * row_spacing
        
        # 左侧单元格
        cell_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(0.2), y, h_w, row_h)
        cell_l.fill.solid()
        cell_l.fill.fore_color.rgb = DARK_GRAY_BG
        cell_l.line.color.rgb = ORANGE
        cell_l.line.dash_style = 4 # MSO_LINE.DASH
        p = cell_l.text_frame.paragraphs[0]
        p.text = "  " + ai_field
        p.font.color.rgb = ORANGE
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        
        # 警告图标
        icon_l = slide.shapes.add_textbox(table_left + Inches(0.2) + h_w - Inches(0.3), y - Inches(0.02), Inches(0.3), row_h)
        p_icon = icon_l.text_frame.paragraphs[0]
        p_icon.text = "⚠"
        p_icon.font.color.rgb = ORANGE
        p_icon.font.size = Pt(10)
        p_icon.alignment = PP_ALIGN.RIGHT

        # 右侧单元格
        cell_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(2.9), y, h_w, row_h)
        cell_r.fill.solid()
        cell_r.fill.fore_color.rgb = DARK_GRAY_BG
        cell_r.line.color.rgb = CHART_GRAY
        p = cell_r.text_frame.paragraphs[0]
        p.text = "  " + sap_field
        p.font.color.rgb = WHITE
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        
        # 勾选图标
        icon_r = slide.shapes.add_textbox(table_left + Inches(2.9) + h_w - Inches(0.3), y - Inches(0.02), Inches(0.3), row_h)
        p_icon = icon_r.text_frame.paragraphs[0]
        p_icon.text = "✔"
        p_icon.font.color.rgb = GRAY_TEXT
        p_icon.font.size = Pt(10)
        p_icon.alignment = PP_ALIGN.RIGHT

    # 4. 右侧：信息面板区域
    def add_info_panel(y_pos, icon_char, title_text, desc_text, highlights=[]):
        box_w = Inches(6.2)
        box_h = Inches(1.1)
        box_x = Inches(6.5)
        
        # 面板背景
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, y_pos, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_BG
        box.line.color.rgb = PANEL_BORDER
        
        # 图标
        icon_box = slide.shapes.add_textbox(box_x + Inches(0.1), y_pos + Inches(0.15), Inches(0.8), Inches(0.8))
        p = icon_box.text_frame.paragraphs[0]
        p.text = icon_char
        p.font.size = Pt(28)
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title = slide.shapes.add_textbox(box_x + Inches(1.0), y_pos + Inches(0.1), box_w - Inches(1.1), Inches(0.3))
        p = title.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 描述文本
        desc = slide.shapes.add_textbox(box_x + Inches(1.0), y_pos + Inches(0.4), box_w - Inches(1.2), Inches(0.6))
        desc.text_frame.word_wrap = True
        p = desc.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        
        if not highlights:
            p.text = desc_text
            p.font.color.rgb = WHITE
        else:
            import re
            pattern = '(' + '|'.join(map(re.escape, highlights)) + ')'
            parts = re.split(pattern, desc_text)
            for part in parts:
                if part in highlights:
                    run = p.add_run()
                    run.text = part
                    run.font.color.rgb = ORANGE
                    run.font.bold = True
                elif part:
                    run = p.add_run()
                    run.text = part
                    run.font.color.rgb = WHITE

    # 添加四个信息面板
    add_info_panel(Inches(1.8), "🤖", "概念误导", "Copilot 完全混淆“配额”与“货源”业务模型，代码完全不可用。")
    add_info_panel(Inches(3.1), "🌩️", "严重幻觉", "Claude 虚构 9 处核心字段（占比 50%），引发 21 个连锁语法错误。", highlights=["9", "50%", "21"])
    add_info_panel(Inches(4.4), "🔒", "接口违规", "AI 无法识别 SAP FM 接口必须使用 DDIC 类型的强制规则。")
    add_info_panel(Inches(5.7), "⏱️", "效率停滞", "修正幻觉字段与重写逻辑的成本已抵消 AI 生成的便利。")

    # 5. 页码
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.4))
    p = page_num.text_frame.paragraphs[0]
    p.text = "2/4"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT



# ── Slide 3 ──

def build_slide_3(slide):
    # 背景设置 (深色主题)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x14, 0x16, 0x1A)
    bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    
    run1 = p.add_run()
    run1.text = "高复杂度：效率反降 "
    run1.font.size = Pt(32)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run1.font.name = FONT_NAME

    run2 = p.add_run()
    run2.text = "60%"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35) # 红色
    run2.font.name = FONT_NAME

    run3 = p.add_run()
    run3.text = " 的修复陷阱"
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run3.font.name = FONT_NAME

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "案例C — 跨工厂 STO 报表（18+ 关联表）"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p_sub.font.name = FONT_NAME

    # 左上区域：-60% 效率损失 (使用文本和箭头代替复杂插图)
    loss_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.2), Inches(2.5), Inches(1.5))
    tf_loss = loss_box.text_frame
    p_loss1 = tf_loss.paragraphs[0]
    p_loss1.text = "-60%"
    p_loss1.font.size = Pt(48)
    p_loss1.font.bold = True
    p_loss1.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_loss1.alignment = PP_ALIGN.CENTER

    p_loss2 = tf_loss.add_paragraph()
    p_loss2.text = "效率损失"
    p_loss2.font.size = Pt(24)
    p_loss2.font.bold = True
    p_loss2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_loss2.alignment = PP_ALIGN.CENTER

    # 向下箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.8), Inches(3.4), Inches(0.8), Inches(1.0))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE5, 0x39, 0x35)
    arrow.line.fill.background()
    arrow.rotation = -45

    # 右上区域：条形图对比
    box_tr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.3), Inches(2.2))
    box_tr.fill.solid()
    box_tr.fill.fore_color.rgb = RGBColor(0x2A, 0x2D, 0x35)
    box_tr.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # AI 修复标签
    lbl_ai = slide.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(2.0), Inches(0.4))
    lbl_ai.text_frame.text = "AI 修复"
    lbl_ai.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lbl_ai.text_frame.paragraphs[0].font.size = Pt(14)

    # AI 修复条 (红色)
    bar_ai = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.4), Inches(4.7), Inches(0.4))
    bar_ai.fill.solid()
    bar_ai.fill.fore_color.rgb = RGBColor(0x9E, 0x1B, 0x22)
    bar_ai.line.fill.background()
    
    txt_ai = slide.shapes.add_textbox(Inches(9.5), Inches(2.4), Inches(2.8), Inches(0.4))
    txt_ai.text_frame.text = "8 人天 (远超手写)"
    txt_ai.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txt_ai.text_frame.paragraphs[0].font.size = Pt(12)
    txt_ai.text_frame.paragraphs[0].font.bold = True
    txt_ai.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 手写开发标签
    lbl_man = slide.shapes.add_textbox(Inches(7.7), Inches(2.9), Inches(2.0), Inches(0.4))
    lbl_man.text_frame.text = "手写开发"
    lbl_man.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lbl_man.text_frame.paragraphs[0].font.size = Pt(14)

    # 手写开发条 (蓝灰色)
    bar_man = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(3.3), Inches(3.0), Inches(0.4))
    bar_man.fill.solid()
    bar_man.fill.fore_color.rgb = RGBColor(0x4A, 0x62, 0x78)
    bar_man.line.fill.background()
    
    txt_man = slide.shapes.add_textbox(Inches(9.0), Inches(3.3), Inches(1.6), Inches(0.4))
    txt_man.text_frame.text = "5 人天"
    txt_man.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txt_man.text_frame.paragraphs[0].font.size = Pt(12)
    txt_man.text_frame.paragraphs[0].font.bold = True
    txt_man.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 左下区域：折线图与警告
    box_bl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(6.5), Inches(2.8))
    box_bl.fill.solid()
    box_bl.fill.fore_color.rgb = RGBColor(0x1A, 0x1C, 0x22)
    box_bl.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # 坐标轴
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(4.6), Inches(1.0), Inches(6.6)).line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.6), Inches(4.0), Inches(6.6)).line.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # 坐标轴标签
    y_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(0.4), Inches(1.5))
    y_label.text_frame.word_wrap = True
    y_label.text_frame.text = "错\n误\n数\n量"
    y_label.text_frame.paragraphs[0].font.size = Pt(10)
    y_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    y_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    x_label = slide.shapes.add_textbox(Inches(3.2), Inches(6.7), Inches(1.0), Inches(0.4))
    x_label.text_frame.text = "修复迭代"
    x_label.text_frame.paragraphs[0].font.size = Pt(10)
    x_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # 数据点与连线 (3 -> 2 -> 9)
    pt1_x, pt1_y = Inches(1.4), Inches(6.1)
    pt2_x, pt2_y = Inches(2.3), Inches(6.3)
    pt3_x, pt3_y = Inches(3.5), Inches(5.0)

    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt1_x, pt1_y, pt2_x, pt2_y)
    line1.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    line1.line.width = Pt(2)
    
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt2_x, pt2_y, pt3_x, pt3_y)
    line2.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    line2.line.width = Pt(2)

    for px, py in [(pt1_x, pt1_y), (pt2_x, pt2_y), (pt3_x, pt3_y)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.05), py - Inches(0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        dot.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
        dot.line.width = Pt(1.5)

    # 数据点标签
    lbl_pt1 = slide.shapes.add_textbox(pt1_x - Inches(0.2), pt1_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt1.text_frame.text = "3"
    lbl_pt1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    lbl_pt2 = slide.shapes.add_textbox(pt2_x - Inches(0.2), pt2_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt2.text_frame.text = "2"
    lbl_pt2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    lbl_pt3 = slide.shapes.add_textbox(pt3_x - Inches(0.1), pt3_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt3.text_frame.text = "9"
    lbl_pt3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 趋势标注
    anno = slide.shapes.add_textbox(Inches(2.4), Inches(5.4), Inches(1.0), Inches(0.3))
    anno.text_frame.text = "3→2→9"
    anno.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    anno.text_frame.paragraphs[0].font.size = Pt(10)
    anno.rotation = -35

    # 警告区域 (修复爆炸)
    warn_icon = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.2), Inches(4.8), Inches(0.3), Inches(0.3))
    warn_icon.fill.solid()
    warn_icon.fill.fore_color.rgb = RGBColor(0xE5, 0x39, 0x35)
    warn_icon.line.fill.background()
    
    warn_ex = slide.shapes.add_textbox(Inches(4.15), Inches(4.8), Inches(0.4), Inches(0.3))
    warn_ex.text_frame.text = "!"
    warn_ex.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    warn_ex.text_frame.paragraphs[0].font.bold = True
    warn_ex.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    warn_title = slide.shapes.add_textbox(Inches(4.6), Inches(4.75), Inches(2.0), Inches(0.4))
    warn_title.text_frame.text = "修复爆炸"
    warn_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    warn_title.text_frame.paragraphs[0].font.size = Pt(16)
    warn_title.text_frame.paragraphs[0].font.bold = True

    warn_desc = slide.shapes.add_textbox(Inches(4.2), Inches(5.3), Inches(2.6), Inches(1.5))
    tf_wd = warn_desc.text_frame
    tf_wd.word_wrap = True
    p_wd = tf_wd.paragraphs[0]
    p_wd.font.size = Pt(12)
    p_wd.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    
    run_wd1 = p_wd.add_run()
    run_wd1.text = "错误呈 "
    run_wd2 = p_wd.add_run()
    run_wd2.text = "3→2→9"
    run_wd2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    run_wd3 = p_wd.add_run()
    run_wd3.text = " 非线性增长，前序修正引发后续逻辑大规模崩溃"

    # 右下区域：要点列表
    # 要点 1: 解析能力
    icon1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(4.5), Inches(0.8), Inches(0.8))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    icon1.line.fill.background()
    
    title1 = slide.shapes.add_textbox(Inches(8.5), Inches(4.4), Inches(4.0), Inches(0.4))
    title1.text_frame.text = "解析能力"
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title1.text_frame.paragraphs[0].font.size = Pt(16)
    title1.text_frame.paragraphs[0].font.bold = True

    desc1 = slide.shapes.add_textbox(Inches(8.5), Inches(4.8), Inches(4.5), Inches(0.8))
    desc1.text_frame.word_wrap = True
    p_d1 = desc1.text_frame.paragraphs[0]
    p_d1.font.size = Pt(12)
    p_d1.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p_d1.add_run().text = "Claude 虽支持文档解析，但在处理 "
    r_d1 = p_d1.add_run()
    r_d1.text = "10+"
    r_d1.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_d1.add_run().text = " 项虚构数据字典时表现乏力"

    # 分隔线
    sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.5), Inches(5.7), Inches(12.8), Inches(5.7))
    sep.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    sep.line.dash_style = 2

    # 要点 2: 实测评价
    icon2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(6.0), Inches(0.8), Inches(0.8))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    icon2.line.fill.background()

    title2 = slide.shapes.add_textbox(Inches(8.5), Inches(5.9), Inches(4.0), Inches(0.4))
    title2.text_frame.text = "实测评价"
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title2.text_frame.paragraphs[0].font.size = Pt(16)
    title2.text_frame.paragraphs[0].font.bold = True

    desc2 = slide.shapes.add_textbox(Inches(8.5), Inches(6.3), Inches(4.5), Inches(0.8))
    desc2.text_frame.word_wrap = True
    p_d2 = desc2.text_frame.paragraphs[0]
    p_d2.font.size = Pt(12)
    p_d2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p_d2.text = "底层数据结构理解缺失，导致“修复成本 > 直接重写”"

    # 页脚
    footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    footer.text_frame.text = "第 3 页"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT



# ── Slide 4 ──

def build_slide_4(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(0x0A, 0x14, 0x28)
    GOLD = RGBColor(0xFF, 0xDF, 0x8C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
    DARK_BLUE_BG = RGBColor(0x15, 0x28, 0x42)
    BORDER_BLUE = RGBColor(0x3A, 0x55, 0x75)
    FONT_NAME = "Microsoft YaHei"

    # Set Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 的应用边界"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = FONT_NAME

    # Header Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效能总览"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # --- Left Section (Quadrant Chart) ---
    # Shadow/Offset Border
    left_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(6.5), Inches(5.5))
    left_shadow.fill.background()
    left_shadow.line.color.rgb = GOLD
    left_shadow.line.width = Pt(1.5)

    # Main Left Background
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = DARK_BLUE_BG
    left_bg.line.color.rgb = BORDER_BLUE
    left_bg.line.width = Pt(1)

    # Axes
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(1.0), Inches(1.8))
    y_axis.line.color.rgb = GOLD
    y_axis.line.width = Pt(2)
    y_axis.line.end_arrowhead = 2

    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(6.8), Inches(6.8))
    x_axis.line.color.rgb = GOLD
    x_axis.line.width = Pt(2)
    x_axis.line.end_arrowhead = 2

    # Axis Labels
    y_label = slide.shapes.add_textbox(Inches(0.2), Inches(4.0), Inches(1.5), Inches(0.5))
    y_label.rotation = -90
    y_p = y_label.text_frame.paragraphs[0]
    y_p.text = "提效百分比"
    y_p.font.size = Pt(14)
    y_p.font.color.rgb = GOLD
    y_p.font.name = FONT_NAME

    x_label = slide.shapes.add_textbox(Inches(3.5), Inches(6.9), Inches(2.0), Inches(0.5))
    x_p = x_label.text_frame.paragraphs[0]
    x_p.text = "任务复杂度"
    x_p.font.size = Pt(14)
    x_p.font.color.rgb = GOLD
    x_p.font.name = FONT_NAME

    # Quadrants Helper Function
    def add_quadrant(left, top, width, height, title, icon_text, main_text, sub_text, is_highlight=False):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1C, 0x33, 0x50) if not is_highlight else RGBColor(0x25, 0x3A, 0x45)
        if is_highlight:
            box.line.color.rgb = GOLD
            box.line.width = Pt(1.5)
        else:
            box.line.color.rgb = BORDER_BLUE
            box.line.width = Pt(1)

        # Title
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Icon
        icon_box = slide.shapes.add_textbox(left + width - Inches(0.6), top + Inches(0.1), Inches(0.5), Inches(0.5))
        p_icon = icon_box.text_frame.paragraphs[0]
        p_icon.text = icon_text
        p_icon.font.size = Pt(18)
        p_icon.font.color.rgb = GOLD if is_highlight else LIGHT_GRAY
        p_icon.font.name = "Segoe UI Symbol"

        # Main Text
        if main_text:
            main_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(0.6))
            main_p = main_box.text_frame.paragraphs[0]
            main_p.text = main_text
            main_p.font.size = Pt(26)
            main_p.font.bold = True
            main_p.font.color.rgb = GOLD
            main_p.alignment = PP_ALIGN.CENTER
            main_p.font.name = FONT_NAME

        # Sub Text
        sub_box = slide.shapes.add_textbox(left, top + Inches(1.3), width, Inches(0.4))
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = sub_text
        sub_p.font.size = Pt(13)
        sub_p.font.color.rgb = LIGHT_GRAY
        sub_p.alignment = PP_ALIGN.CENTER
        sub_p.font.name = FONT_NAME

        # Special handling for Bottom-Right broken link icon
        if title == "应用禁区":
            link_icon = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.4), Inches(1.1), Inches(1.0))
            link_p = link_icon.text_frame.paragraphs[0]
            link_p.text = "🔗"
            link_p.font.size = Pt(45)
            link_p.font.color.rgb = LIGHT_GRAY
            link_p.alignment = PP_ALIGN.CENTER
            link_p.font.name = "Segoe UI Emoji"

            strike = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(1.0), top + Inches(0.5), left + Inches(1.7), top + Inches(1.2))
            strike.line.color.rgb = GOLD
            strike.line.width = Pt(3)

    # Add 4 Quadrants
    add_quadrant(Inches(1.2), Inches(2.0), Inches(2.6), Inches(2.2), "简单场景", "☑", "50% 提效", "辅助逻辑，高收益")
    add_quadrant(Inches(1.2), Inches(4.4), Inches(2.6), Inches(2.2), "中等场景", "⊖", "无提升", "效率持平")
    add_quadrant(Inches(4.0), Inches(2.0), Inches(2.6), Inches(2.2), "复杂场景", "⚠", "工作量激增", "效率倒挂，风险高", is_highlight=True)
    add_quadrant(Inches(4.0), Inches(4.4), Inches(2.6), Inches(2.2), "应用禁区", "⊘", "", "不可靠，需人工完全接管")

    # Trend Curves
    curve_tl = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(1.5), Inches(4.0), Inches(3.6), Inches(3.0))
    curve_tl.line.color.rgb = GOLD
    curve_tl.line.width = Pt(2)
    curve_tl.line.end_arrowhead = 2

    curve_tr = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(4.2), Inches(2.5), Inches(6.4), Inches(4.0))
    curve_tr.line.color.rgb = GOLD
    curve_tr.line.width = Pt(2)
    curve_tr.line.end_arrowhead = 2

    # --- Right Section ---

    # 1. Bottlenecks Box
    bot_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.3), Inches(1.8))
    bot_bg.fill.solid()
    bot_bg.fill.fore_color.rgb = DARK_BLUE_BG
    bot_bg.line.color.rgb = BORDER_BLUE
    bot_bg.line.width = Pt(1)

    bot_title = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(5.3), Inches(0.4))
    bot_p = bot_title.text_frame.paragraphs[0]
    bot_p.text = "四大瓶颈"
    bot_p.font.size = Pt(16)
    bot_p.font.bold = True
    bot_p.font.color.rgb = GOLD
    bot_p.alignment = PP_ALIGN.CENTER
    bot_p.font.name = FONT_NAME

    def add_list_item(left, top, icon, text):
        tb = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{icon}  {text}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(2.2), "☁", "数据字典幻觉")
    add_list_item(Inches(10.3), Inches(2.2), "🔗", "接口规则缺失")
    add_list_item(Inches(7.8), Inches(2.7), "📈", "修复成本反超")
    add_list_item(Inches(10.3), Inches(2.7), "💬", "提示词被忽略")

    # 2. Tools Box
    tool_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.6), Inches(5.3), Inches(1.8))
    tool_bg.fill.solid()
    tool_bg.fill.fore_color.rgb = DARK_BLUE_BG
    tool_bg.line.color.rgb = BORDER_BLUE
    tool_bg.line.width = Pt(1)

    tool_title = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(5.3), Inches(0.4))
    tool_p = tool_title.text_frame.paragraphs[0]
    tool_p.text = "工具选型"
    tool_p.font.size = Pt(16)
    tool_p.font.bold = True
    tool_p.font.color.rgb = GOLD
    tool_p.alignment = PP_ALIGN.CENTER
    tool_p.font.name = FONT_NAME

    # Claude Code
    cc_title = slide.shapes.add_textbox(Inches(7.8), Inches(4.1), Inches(2.5), Inches(0.4))
    cc_p = cc_title.text_frame.paragraphs[0]
    cc_p.text = "⚡ Claude Code"
    cc_p.font.size = Pt(15)
    cc_p.font.bold = True
    cc_p.font.color.rgb = WHITE
    cc_p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(4.5), "☑", "业务理解优秀")
    add_list_item(Inches(7.8), Inches(4.9), "☑", "多文件解析强")

    # Copilot
    cp_title = slide.shapes.add_textbox(Inches(10.3), Inches(4.1), Inches(2.5), Inches(0.4))
    cp_p = cp_title.text_frame.paragraphs[0]
    cp_p.text = "🤖 Copilot"
    cp_p.font.size = Pt(15)
    cp_p.font.bold = True
    cp_p.font.color.rgb = WHITE
    cp_p.font.name = FONT_NAME

    add_list_item(Inches(10.3), Inches(4.5), "☒", "业务理解受限")
    add_list_item(Inches(10.3), Inches(4.9), "☒", "多文件解析弱")

    # 3. Conclusion Box
    conc_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(5.7), Inches(5.3), Inches(1.4))
    conc_shadow.fill.background()
    conc_shadow.line.color.rgb = GOLD
    conc_shadow.line.width = Pt(1.5)

    conc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.6), Inches(5.3), Inches(1.4))
    conc_bg.fill.solid()
    conc_bg.fill.fore_color.rgb = RGBColor(0x1C, 0x2A, 0x35)
    conc_bg.line.color.rgb = BORDER_BLUE
    conc_bg.line.width = Pt(1)

    rb_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.8), Inches(4.9), Inches(1.1))
    rb_tf = rb_text.text_frame
    rb_tf.word_wrap = True
    rb_p = rb_tf.paragraphs[0]

    run1 = rb_p.add_run()
    run1.text = "核心建议："
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = GOLD
    run1.font.name = FONT_NAME

    run2 = rb_p.add_run()
    run2.text = "当前 AI 尚未具备处理 ABAP 中高难度任务的能力，仅推荐用于"
    run2.font.size = Pt(15)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME

    run3 = rb_p.add_run()
    run3.text = "辅助简单逻辑"
    run3.font.size = Pt(15)
    run3.font.bold = True
    run3.font.color.rgb = GOLD
    run3.font.name = FONT_NAME

    run4 = rb_p.add_run()
    run4.text = "。"
    run4.font.size = Pt(15)
    run4.font.color.rgb = WHITE
    run4.font.name = FONT_NAME



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
s1 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_2(s1)
s2 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_3(s2)
s3 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_4(s3)
prs.save(OUTPUT_PATH)

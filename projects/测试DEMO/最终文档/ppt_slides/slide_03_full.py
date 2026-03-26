from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"projects/测试DEMO/最终文档/ppt_slides/slide_03.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 背景设置 (深色主题)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x14, 0x16, 0x1A)
    bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    
    run1 = p.add_run()
    run1.text = "高复杂度：效率反降 "
    run1.font.size = Pt(32)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run1.font.name = FONT_NAME

    run2 = p.add_run()
    run2.text = "60%"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35) # 红色
    run2.font.name = FONT_NAME

    run3 = p.add_run()
    run3.text = " 的修复陷阱"
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run3.font.name = FONT_NAME

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "案例C — 跨工厂 STO 报表（18+ 关联表）"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p_sub.font.name = FONT_NAME

    # 左上区域：-60% 效率损失 (使用文本和箭头代替复杂插图)
    loss_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.2), Inches(2.5), Inches(1.5))
    tf_loss = loss_box.text_frame
    p_loss1 = tf_loss.paragraphs[0]
    p_loss1.text = "-60%"
    p_loss1.font.size = Pt(48)
    p_loss1.font.bold = True
    p_loss1.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_loss1.alignment = PP_ALIGN.CENTER

    p_loss2 = tf_loss.add_paragraph()
    p_loss2.text = "效率损失"
    p_loss2.font.size = Pt(24)
    p_loss2.font.bold = True
    p_loss2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_loss2.alignment = PP_ALIGN.CENTER

    # 向下箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.8), Inches(3.4), Inches(0.8), Inches(1.0))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE5, 0x39, 0x35)
    arrow.line.fill.background()
    arrow.rotation = -45

    # 右上区域：条形图对比
    box_tr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.3), Inches(2.2))
    box_tr.fill.solid()
    box_tr.fill.fore_color.rgb = RGBColor(0x2A, 0x2D, 0x35)
    box_tr.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # AI 修复标签
    lbl_ai = slide.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(2.0), Inches(0.4))
    lbl_ai.text_frame.text = "AI 修复"
    lbl_ai.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lbl_ai.text_frame.paragraphs[0].font.size = Pt(14)

    # AI 修复条 (红色)
    bar_ai = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.4), Inches(4.7), Inches(0.4))
    bar_ai.fill.solid()
    bar_ai.fill.fore_color.rgb = RGBColor(0x9E, 0x1B, 0x22)
    bar_ai.line.fill.background()
    
    txt_ai = slide.shapes.add_textbox(Inches(9.5), Inches(2.4), Inches(2.8), Inches(0.4))
    txt_ai.text_frame.text = "8 人天 (远超手写)"
    txt_ai.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txt_ai.text_frame.paragraphs[0].font.size = Pt(12)
    txt_ai.text_frame.paragraphs[0].font.bold = True
    txt_ai.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 手写开发标签
    lbl_man = slide.shapes.add_textbox(Inches(7.7), Inches(2.9), Inches(2.0), Inches(0.4))
    lbl_man.text_frame.text = "手写开发"
    lbl_man.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lbl_man.text_frame.paragraphs[0].font.size = Pt(14)

    # 手写开发条 (蓝灰色)
    bar_man = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(3.3), Inches(3.0), Inches(0.4))
    bar_man.fill.solid()
    bar_man.fill.fore_color.rgb = RGBColor(0x4A, 0x62, 0x78)
    bar_man.line.fill.background()
    
    txt_man = slide.shapes.add_textbox(Inches(9.0), Inches(3.3), Inches(1.6), Inches(0.4))
    txt_man.text_frame.text = "5 人天"
    txt_man.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txt_man.text_frame.paragraphs[0].font.size = Pt(12)
    txt_man.text_frame.paragraphs[0].font.bold = True
    txt_man.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 左下区域：折线图与警告
    box_bl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(6.5), Inches(2.8))
    box_bl.fill.solid()
    box_bl.fill.fore_color.rgb = RGBColor(0x1A, 0x1C, 0x22)
    box_bl.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # 坐标轴
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(4.6), Inches(1.0), Inches(6.6)).line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.6), Inches(4.0), Inches(6.6)).line.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # 坐标轴标签
    y_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(0.4), Inches(1.5))
    y_label.text_frame.word_wrap = True
    y_label.text_frame.text = "错\n误\n数\n量"
    y_label.text_frame.paragraphs[0].font.size = Pt(10)
    y_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    y_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    x_label = slide.shapes.add_textbox(Inches(3.2), Inches(6.7), Inches(1.0), Inches(0.4))
    x_label.text_frame.text = "修复迭代"
    x_label.text_frame.paragraphs[0].font.size = Pt(10)
    x_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # 数据点与连线 (3 -> 2 -> 9)
    pt1_x, pt1_y = Inches(1.4), Inches(6.1)
    pt2_x, pt2_y = Inches(2.3), Inches(6.3)
    pt3_x, pt3_y = Inches(3.5), Inches(5.0)

    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt1_x, pt1_y, pt2_x, pt2_y)
    line1.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    line1.line.width = Pt(2)
    
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt2_x, pt2_y, pt3_x, pt3_y)
    line2.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    line2.line.width = Pt(2)

    for px, py in [(pt1_x, pt1_y), (pt2_x, pt2_y), (pt3_x, pt3_y)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.05), py - Inches(0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        dot.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
        dot.line.width = Pt(1.5)

    # 数据点标签
    lbl_pt1 = slide.shapes.add_textbox(pt1_x - Inches(0.2), pt1_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt1.text_frame.text = "3"
    lbl_pt1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    lbl_pt2 = slide.shapes.add_textbox(pt2_x - Inches(0.2), pt2_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt2.text_frame.text = "2"
    lbl_pt2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    lbl_pt3 = slide.shapes.add_textbox(pt3_x - Inches(0.1), pt3_y - Inches(0.3), Inches(0.4), Inches(0.3))
    lbl_pt3.text_frame.text = "9"
    lbl_pt3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 趋势标注
    anno = slide.shapes.add_textbox(Inches(2.4), Inches(5.4), Inches(1.0), Inches(0.3))
    anno.text_frame.text = "3→2→9"
    anno.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    anno.text_frame.paragraphs[0].font.size = Pt(10)
    anno.rotation = -35

    # 警告区域 (修复爆炸)
    warn_icon = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.2), Inches(4.8), Inches(0.3), Inches(0.3))
    warn_icon.fill.solid()
    warn_icon.fill.fore_color.rgb = RGBColor(0xE5, 0x39, 0x35)
    warn_icon.line.fill.background()
    
    warn_ex = slide.shapes.add_textbox(Inches(4.15), Inches(4.8), Inches(0.4), Inches(0.3))
    warn_ex.text_frame.text = "!"
    warn_ex.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    warn_ex.text_frame.paragraphs[0].font.bold = True
    warn_ex.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    warn_title = slide.shapes.add_textbox(Inches(4.6), Inches(4.75), Inches(2.0), Inches(0.4))
    warn_title.text_frame.text = "修复爆炸"
    warn_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    warn_title.text_frame.paragraphs[0].font.size = Pt(16)
    warn_title.text_frame.paragraphs[0].font.bold = True

    warn_desc = slide.shapes.add_textbox(Inches(4.2), Inches(5.3), Inches(2.6), Inches(1.5))
    tf_wd = warn_desc.text_frame
    tf_wd.word_wrap = True
    p_wd = tf_wd.paragraphs[0]
    p_wd.font.size = Pt(12)
    p_wd.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    
    run_wd1 = p_wd.add_run()
    run_wd1.text = "错误呈 "
    run_wd2 = p_wd.add_run()
    run_wd2.text = "3→2→9"
    run_wd2.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    run_wd3 = p_wd.add_run()
    run_wd3.text = " 非线性增长，前序修正引发后续逻辑大规模崩溃"

    # 右下区域：要点列表
    # 要点 1: 解析能力
    icon1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(4.5), Inches(0.8), Inches(0.8))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    icon1.line.fill.background()
    
    title1 = slide.shapes.add_textbox(Inches(8.5), Inches(4.4), Inches(4.0), Inches(0.4))
    title1.text_frame.text = "解析能力"
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title1.text_frame.paragraphs[0].font.size = Pt(16)
    title1.text_frame.paragraphs[0].font.bold = True

    desc1 = slide.shapes.add_textbox(Inches(8.5), Inches(4.8), Inches(4.5), Inches(0.8))
    desc1.text_frame.word_wrap = True
    p_d1 = desc1.text_frame.paragraphs[0]
    p_d1.font.size = Pt(12)
    p_d1.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p_d1.add_run().text = "Claude 虽支持文档解析，但在处理 "
    r_d1 = p_d1.add_run()
    r_d1.text = "10+"
    r_d1.font.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    p_d1.add_run().text = " 项虚构数据字典时表现乏力"

    # 分隔线
    sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.5), Inches(5.7), Inches(12.8), Inches(5.7))
    sep.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    sep.line.dash_style = 2

    # 要点 2: 实测评价
    icon2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(6.0), Inches(0.8), Inches(0.8))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    icon2.line.fill.background()

    title2 = slide.shapes.add_textbox(Inches(8.5), Inches(5.9), Inches(4.0), Inches(0.4))
    title2.text_frame.text = "实测评价"
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title2.text_frame.paragraphs[0].font.size = Pt(16)
    title2.text_frame.paragraphs[0].font.bold = True

    desc2 = slide.shapes.add_textbox(Inches(8.5), Inches(6.3), Inches(4.5), Inches(0.8))
    desc2.text_frame.word_wrap = True
    p_d2 = desc2.text_frame.paragraphs[0]
    p_d2.font.size = Pt(12)
    p_d2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p_d2.text = "底层数据结构理解缺失，导致“修复成本 > 直接重写”"

    # 页脚
    footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    footer.text_frame.text = "第 3 页"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"projects/测试DEMO/最终文档/ppt_slides/slide_01.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BLUE_DARK = RGBColor(0x1F, 0x4E, 0x79)
    BLUE_PRIMARY = RGBColor(0x00, 0x70, 0xC0)
    BLUE_LIGHT_BG = RGBColor(0xF4, 0xF8, 0xFC)
    BLUE_BORDER = RGBColor(0xBD, 0xD7, 0xEE)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    GRAY_BAR = RGBColor(0xA6, 0xA6, 0xA6)
    ORANGE_WARN = RGBColor(0xED, 0x7D, 0x31)
    BLACK_TEXT = RGBColor(0x33, 0x33, 0x33)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "简单场景：AI 初显 50% 提效潜力 ⚡"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_NAME

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "案例A —— 标准 ALV 销售订单报表验证"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = FONT_NAME

    # 3. Page Indicator (Top Right)
    page_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.8), Inches(0.6), Inches(1.0), Inches(0.4))
    page_shape.fill.solid()
    page_shape.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    page_shape.line.fill.background()
    tf = page_shape.text_frame
    tf.text = "P1 / 4"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = BLUE_PRIMARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 4. Left Column - Core Points Heading
    heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3), Inches(0.5))
    tf = heading_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心要点"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = FONT_NAME

    # 5. Left Column - Core Points Items
    items = [
        ("⚡", "1. 效率突破: ", "Claude Code 仅需 30 分钟完成开发，较手写提效 50%。"),
        ("🤖", "2. 工具差距: ", "GitHub Copilot 耗时与手写持平，且因 SQL 不兼容导致运行崩溃。"),
        ("🔀", "3. 交互对比: ", "Claude 仅需 2 轮提示即可运行，Copilot 需 5 轮以上人工干预。"),
        ("🔍", "4. 核心瓶颈: ", "[地址取数逻辑] 与 [过账状态字段] 是 AI 初始生成的共同盲区。")
    ]

    y_offset = 2.7
    for i, (icon, label, desc) in enumerate(items):
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(0.6), Inches(0.6))
        icon_box.text_frame.text = icon
        icon_box.text_frame.paragraphs[0].font.size = Pt(24)

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.2), y_offset, Inches(4.8), Inches(0.8))
        text_box.text_frame.word_wrap = True
        p = text_box.text_frame.paragraphs[0]
        
        run1 = p.add_run()
        run1.text = label
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = BLUE_DARK
        run1.font.name = FONT_NAME
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(14)
        run2.font.color.rgb = BLACK_TEXT
        run2.font.name = FONT_NAME

        # Separator line
        if i < len(items) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), y_offset + 0.85, Inches(6.0), y_offset + 0.85)
            line.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        
        y_offset += 1.0

    # 6. Right Column - Chart Box Background
    chart_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.3), Inches(4.2))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = BLUE_LIGHT_BG
    chart_bg.line.color.rgb = BLUE_BORDER
    chart_bg.line.width = Pt(1.5)

    # Chart Title
    chart_title = slide.shapes.add_textbox(Inches(6.7), Inches(2.0), Inches(4), Inches(0.5))
    tf = chart_title.text_frame
    p = tf.paragraphs[0]
    p.text = "开发耗时对比 (分钟)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = FONT_NAME

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ["手写/Copilot\n(Manual/AI)", "Claude Code\n(AI)"]
    chart_data.add_series('耗时', [60, 30])

    x, y, cx, cy = Inches(6.6), Inches(2.8), Inches(5.8), Inches(2.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.plots[0].has_data_labels = False # We will add custom text boxes for labels
    chart.plots[0].gap_width = 100

    # Customize bar colors
    series = chart.plots[0].series[0]
    pt0 = series.points[0]
    pt0.format.fill.solid()
    pt0.format.fill.fore_color.rgb = GRAY_BAR
    pt1 = series.points[1]
    pt1.format.fill.solid()
    pt1.format.fill.fore_color.rgb = BLUE_PRIMARY

    # Custom Data Labels
    lbl1 = slide.shapes.add_textbox(Inches(11.9), Inches(4.15), Inches(1), Inches(0.4))
    lbl1.text_frame.text = "⚠️ 60"
    lbl1.text_frame.paragraphs[0].font.size = Pt(14)
    lbl1.text_frame.paragraphs[0].font.name = FONT_NAME

    lbl2 = slide.shapes.add_textbox(Inches(9.6), Inches(3.25), Inches(1), Inches(0.4))
    lbl2.text_frame.text = "⚡ 30"
    lbl2.text_frame.paragraphs[0].font.size = Pt(14)
    lbl2.text_frame.paragraphs[0].font.name = FONT_NAME

    # Chart Callout (50% 提效)
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, Inches(9.1), Inches(2.4), Inches(1.6), Inches(0.6))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    callout.line.color.rgb = BLUE_BORDER
    tf = callout.text_frame
    p = tf.paragraphs[0]
    p.text = "50% 提效"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Chart Warning Text
    warn_text = slide.shapes.add_textbox(Inches(10.3), Inches(4.6), Inches(2.5), Inches(0.4))
    tf = warn_text.text_frame
    p = tf.paragraphs[0]
    p.text = "SQL 不兼容/运行崩溃"
    p.font.size = Pt(10)
    p.font.color.rgb = ORANGE_WARN
    p.font.name = FONT_NAME

    # 7. Right Column - Conclusion Box
    conc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(6.2), Inches(6.3), Inches(0.8))
    conc_bg.fill.solid()
    conc_bg.fill.fore_color.rgb = BLUE_LIGHT_BG
    conc_bg.line.color.rgb = BLUE_BORDER
    conc_bg.line.width = Pt(1.5)

    conc_text = slide.shapes.add_textbox(Inches(6.5), Inches(6.35), Inches(6.3), Inches(0.5))
    tf = conc_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run1 = p.add_run()
    run1.text = "✔️ ⚡ AI 初显 "
    run1.font.size = Pt(22)
    run1.font.bold = True
    run1.font.color.rgb = BLACK_TEXT
    run1.font.name = FONT_NAME

    run2 = p.add_run()
    run2.text = "50% 提效潜力"
    run2.font.size = Pt(22)
    run2.font.bold = True
    run2.font.color.rgb = BLUE_PRIMARY
    run2.font.name = FONT_NAME



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)

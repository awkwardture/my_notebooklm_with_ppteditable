from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"projects/测试DEMO/最终文档/ppt_slides/slide_04.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(0x0A, 0x14, 0x28)
    GOLD = RGBColor(0xFF, 0xDF, 0x8C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
    DARK_BLUE_BG = RGBColor(0x15, 0x28, 0x42)
    BORDER_BLUE = RGBColor(0x3A, 0x55, 0x75)
    FONT_NAME = "Microsoft YaHei"

    # Set Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 的应用边界"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = FONT_NAME

    # Header Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效能总览"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # --- Left Section (Quadrant Chart) ---
    # Shadow/Offset Border
    left_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(6.5), Inches(5.5))
    left_shadow.fill.background()
    left_shadow.line.color.rgb = GOLD
    left_shadow.line.width = Pt(1.5)

    # Main Left Background
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = DARK_BLUE_BG
    left_bg.line.color.rgb = BORDER_BLUE
    left_bg.line.width = Pt(1)

    # Axes
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(1.0), Inches(1.8))
    y_axis.line.color.rgb = GOLD
    y_axis.line.width = Pt(2)
    y_axis.line.end_arrowhead = 2

    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(6.8), Inches(6.8))
    x_axis.line.color.rgb = GOLD
    x_axis.line.width = Pt(2)
    x_axis.line.end_arrowhead = 2

    # Axis Labels
    y_label = slide.shapes.add_textbox(Inches(0.2), Inches(4.0), Inches(1.5), Inches(0.5))
    y_label.rotation = -90
    y_p = y_label.text_frame.paragraphs[0]
    y_p.text = "提效百分比"
    y_p.font.size = Pt(14)
    y_p.font.color.rgb = GOLD
    y_p.font.name = FONT_NAME

    x_label = slide.shapes.add_textbox(Inches(3.5), Inches(6.9), Inches(2.0), Inches(0.5))
    x_p = x_label.text_frame.paragraphs[0]
    x_p.text = "任务复杂度"
    x_p.font.size = Pt(14)
    x_p.font.color.rgb = GOLD
    x_p.font.name = FONT_NAME

    # Quadrants Helper Function
    def add_quadrant(left, top, width, height, title, icon_text, main_text, sub_text, is_highlight=False):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1C, 0x33, 0x50) if not is_highlight else RGBColor(0x25, 0x3A, 0x45)
        if is_highlight:
            box.line.color.rgb = GOLD
            box.line.width = Pt(1.5)
        else:
            box.line.color.rgb = BORDER_BLUE
            box.line.width = Pt(1)

        # Title
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Icon
        icon_box = slide.shapes.add_textbox(left + width - Inches(0.6), top + Inches(0.1), Inches(0.5), Inches(0.5))
        p_icon = icon_box.text_frame.paragraphs[0]
        p_icon.text = icon_text
        p_icon.font.size = Pt(18)
        p_icon.font.color.rgb = GOLD if is_highlight else LIGHT_GRAY
        p_icon.font.name = "Segoe UI Symbol"

        # Main Text
        if main_text:
            main_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(0.6))
            main_p = main_box.text_frame.paragraphs[0]
            main_p.text = main_text
            main_p.font.size = Pt(26)
            main_p.font.bold = True
            main_p.font.color.rgb = GOLD
            main_p.alignment = PP_ALIGN.CENTER
            main_p.font.name = FONT_NAME

        # Sub Text
        sub_box = slide.shapes.add_textbox(left, top + Inches(1.3), width, Inches(0.4))
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = sub_text
        sub_p.font.size = Pt(13)
        sub_p.font.color.rgb = LIGHT_GRAY
        sub_p.alignment = PP_ALIGN.CENTER
        sub_p.font.name = FONT_NAME

        # Special handling for Bottom-Right broken link icon
        if title == "应用禁区":
            link_icon = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.4), Inches(1.1), Inches(1.0))
            link_p = link_icon.text_frame.paragraphs[0]
            link_p.text = "🔗"
            link_p.font.size = Pt(45)
            link_p.font.color.rgb = LIGHT_GRAY
            link_p.alignment = PP_ALIGN.CENTER
            link_p.font.name = "Segoe UI Emoji"

            strike = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(1.0), top + Inches(0.5), left + Inches(1.7), top + Inches(1.2))
            strike.line.color.rgb = GOLD
            strike.line.width = Pt(3)

    # Add 4 Quadrants
    add_quadrant(Inches(1.2), Inches(2.0), Inches(2.6), Inches(2.2), "简单场景", "☑", "50% 提效", "辅助逻辑，高收益")
    add_quadrant(Inches(1.2), Inches(4.4), Inches(2.6), Inches(2.2), "中等场景", "⊖", "无提升", "效率持平")
    add_quadrant(Inches(4.0), Inches(2.0), Inches(2.6), Inches(2.2), "复杂场景", "⚠", "工作量激增", "效率倒挂，风险高", is_highlight=True)
    add_quadrant(Inches(4.0), Inches(4.4), Inches(2.6), Inches(2.2), "应用禁区", "⊘", "", "不可靠，需人工完全接管")

    # Trend Curves
    curve_tl = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(1.5), Inches(4.0), Inches(3.6), Inches(3.0))
    curve_tl.line.color.rgb = GOLD
    curve_tl.line.width = Pt(2)
    curve_tl.line.end_arrowhead = 2

    curve_tr = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(4.2), Inches(2.5), Inches(6.4), Inches(4.0))
    curve_tr.line.color.rgb = GOLD
    curve_tr.line.width = Pt(2)
    curve_tr.line.end_arrowhead = 2

    # --- Right Section ---

    # 1. Bottlenecks Box
    bot_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.3), Inches(1.8))
    bot_bg.fill.solid()
    bot_bg.fill.fore_color.rgb = DARK_BLUE_BG
    bot_bg.line.color.rgb = BORDER_BLUE
    bot_bg.line.width = Pt(1)

    bot_title = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(5.3), Inches(0.4))
    bot_p = bot_title.text_frame.paragraphs[0]
    bot_p.text = "四大瓶颈"
    bot_p.font.size = Pt(16)
    bot_p.font.bold = True
    bot_p.font.color.rgb = GOLD
    bot_p.alignment = PP_ALIGN.CENTER
    bot_p.font.name = FONT_NAME

    def add_list_item(left, top, icon, text):
        tb = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{icon}  {text}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(2.2), "☁", "数据字典幻觉")
    add_list_item(Inches(10.3), Inches(2.2), "🔗", "接口规则缺失")
    add_list_item(Inches(7.8), Inches(2.7), "📈", "修复成本反超")
    add_list_item(Inches(10.3), Inches(2.7), "💬", "提示词被忽略")

    # 2. Tools Box
    tool_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.6), Inches(5.3), Inches(1.8))
    tool_bg.fill.solid()
    tool_bg.fill.fore_color.rgb = DARK_BLUE_BG
    tool_bg.line.color.rgb = BORDER_BLUE
    tool_bg.line.width = Pt(1)

    tool_title = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(5.3), Inches(0.4))
    tool_p = tool_title.text_frame.paragraphs[0]
    tool_p.text = "工具选型"
    tool_p.font.size = Pt(16)
    tool_p.font.bold = True
    tool_p.font.color.rgb = GOLD
    tool_p.alignment = PP_ALIGN.CENTER
    tool_p.font.name = FONT_NAME

    # Claude Code
    cc_title = slide.shapes.add_textbox(Inches(7.8), Inches(4.1), Inches(2.5), Inches(0.4))
    cc_p = cc_title.text_frame.paragraphs[0]
    cc_p.text = "⚡ Claude Code"
    cc_p.font.size = Pt(15)
    cc_p.font.bold = True
    cc_p.font.color.rgb = WHITE
    cc_p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(4.5), "☑", "业务理解优秀")
    add_list_item(Inches(7.8), Inches(4.9), "☑", "多文件解析强")

    # Copilot
    cp_title = slide.shapes.add_textbox(Inches(10.3), Inches(4.1), Inches(2.5), Inches(0.4))
    cp_p = cp_title.text_frame.paragraphs[0]
    cp_p.text = "🤖 Copilot"
    cp_p.font.size = Pt(15)
    cp_p.font.bold = True
    cp_p.font.color.rgb = WHITE
    cp_p.font.name = FONT_NAME

    add_list_item(Inches(10.3), Inches(4.5), "☒", "业务理解受限")
    add_list_item(Inches(10.3), Inches(4.9), "☒", "多文件解析弱")

    # 3. Conclusion Box
    conc_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(5.7), Inches(5.3), Inches(1.4))
    conc_shadow.fill.background()
    conc_shadow.line.color.rgb = GOLD
    conc_shadow.line.width = Pt(1.5)

    conc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.6), Inches(5.3), Inches(1.4))
    conc_bg.fill.solid()
    conc_bg.fill.fore_color.rgb = RGBColor(0x1C, 0x2A, 0x35)
    conc_bg.line.color.rgb = BORDER_BLUE
    conc_bg.line.width = Pt(1)

    rb_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.8), Inches(4.9), Inches(1.1))
    rb_tf = rb_text.text_frame
    rb_tf.word_wrap = True
    rb_p = rb_tf.paragraphs[0]

    run1 = rb_p.add_run()
    run1.text = "核心建议："
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = GOLD
    run1.font.name = FONT_NAME

    run2 = rb_p.add_run()
    run2.text = "当前 AI 尚未具备处理 ABAP 中高难度任务的能力，仅推荐用于"
    run2.font.size = Pt(15)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME

    run3 = rb_p.add_run()
    run3.text = "辅助简单逻辑"
    run3.font.size = Pt(15)
    run3.font.bold = True
    run3.font.color.rgb = GOLD
    run3.font.name = FONT_NAME

    run4 = rb_p.add_run()
    run4.text = "。"
    run4.font.size = Pt(15)
    run4.font.color.rgb = WHITE
    run4.font.name = FONT_NAME



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)

def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(0x0A, 0x14, 0x28)
    GOLD = RGBColor(0xFF, 0xDF, 0x8C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
    DARK_BLUE_BG = RGBColor(0x15, 0x28, 0x42)
    BORDER_BLUE = RGBColor(0x3A, 0x55, 0x75)
    FONT_NAME = "Microsoft YaHei"

    # Set Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 的应用边界"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = FONT_NAME

    # Header Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效能总览"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # --- Left Section (Quadrant Chart) ---
    # Shadow/Offset Border
    left_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(6.5), Inches(5.5))
    left_shadow.fill.background()
    left_shadow.line.color.rgb = GOLD
    left_shadow.line.width = Pt(1.5)

    # Main Left Background
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = DARK_BLUE_BG
    left_bg.line.color.rgb = BORDER_BLUE
    left_bg.line.width = Pt(1)

    # Axes
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(1.0), Inches(1.8))
    y_axis.line.color.rgb = GOLD
    y_axis.line.width = Pt(2)
    y_axis.line.end_arrowhead = 2

    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(6.8), Inches(6.8), Inches(6.8))
    x_axis.line.color.rgb = GOLD
    x_axis.line.width = Pt(2)
    x_axis.line.end_arrowhead = 2

    # Axis Labels
    y_label = slide.shapes.add_textbox(Inches(0.2), Inches(4.0), Inches(1.5), Inches(0.5))
    y_label.rotation = -90
    y_p = y_label.text_frame.paragraphs[0]
    y_p.text = "提效百分比"
    y_p.font.size = Pt(14)
    y_p.font.color.rgb = GOLD
    y_p.font.name = FONT_NAME

    x_label = slide.shapes.add_textbox(Inches(3.5), Inches(6.9), Inches(2.0), Inches(0.5))
    x_p = x_label.text_frame.paragraphs[0]
    x_p.text = "任务复杂度"
    x_p.font.size = Pt(14)
    x_p.font.color.rgb = GOLD
    x_p.font.name = FONT_NAME

    # Quadrants Helper Function
    def add_quadrant(left, top, width, height, title, icon_text, main_text, sub_text, is_highlight=False):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1C, 0x33, 0x50) if not is_highlight else RGBColor(0x25, 0x3A, 0x45)
        if is_highlight:
            box.line.color.rgb = GOLD
            box.line.width = Pt(1.5)
        else:
            box.line.color.rgb = BORDER_BLUE
            box.line.width = Pt(1)

        # Title
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Icon
        icon_box = slide.shapes.add_textbox(left + width - Inches(0.6), top + Inches(0.1), Inches(0.5), Inches(0.5))
        p_icon = icon_box.text_frame.paragraphs[0]
        p_icon.text = icon_text
        p_icon.font.size = Pt(18)
        p_icon.font.color.rgb = GOLD if is_highlight else LIGHT_GRAY
        p_icon.font.name = "Segoe UI Symbol"

        # Main Text
        if main_text:
            main_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(0.6))
            main_p = main_box.text_frame.paragraphs[0]
            main_p.text = main_text
            main_p.font.size = Pt(26)
            main_p.font.bold = True
            main_p.font.color.rgb = GOLD
            main_p.alignment = PP_ALIGN.CENTER
            main_p.font.name = FONT_NAME

        # Sub Text
        sub_box = slide.shapes.add_textbox(left, top + Inches(1.3), width, Inches(0.4))
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = sub_text
        sub_p.font.size = Pt(13)
        sub_p.font.color.rgb = LIGHT_GRAY
        sub_p.alignment = PP_ALIGN.CENTER
        sub_p.font.name = FONT_NAME

        # Special handling for Bottom-Right broken link icon
        if title == "应用禁区":
            link_icon = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.4), Inches(1.1), Inches(1.0))
            link_p = link_icon.text_frame.paragraphs[0]
            link_p.text = "🔗"
            link_p.font.size = Pt(45)
            link_p.font.color.rgb = LIGHT_GRAY
            link_p.alignment = PP_ALIGN.CENTER
            link_p.font.name = "Segoe UI Emoji"

            strike = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(1.0), top + Inches(0.5), left + Inches(1.7), top + Inches(1.2))
            strike.line.color.rgb = GOLD
            strike.line.width = Pt(3)

    # Add 4 Quadrants
    add_quadrant(Inches(1.2), Inches(2.0), Inches(2.6), Inches(2.2), "简单场景", "☑", "50% 提效", "辅助逻辑，高收益")
    add_quadrant(Inches(1.2), Inches(4.4), Inches(2.6), Inches(2.2), "中等场景", "⊖", "无提升", "效率持平")
    add_quadrant(Inches(4.0), Inches(2.0), Inches(2.6), Inches(2.2), "复杂场景", "⚠", "工作量激增", "效率倒挂，风险高", is_highlight=True)
    add_quadrant(Inches(4.0), Inches(4.4), Inches(2.6), Inches(2.2), "应用禁区", "⊘", "", "不可靠，需人工完全接管")

    # Trend Curves
    curve_tl = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(1.5), Inches(4.0), Inches(3.6), Inches(3.0))
    curve_tl.line.color.rgb = GOLD
    curve_tl.line.width = Pt(2)
    curve_tl.line.end_arrowhead = 2

    curve_tr = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(4.2), Inches(2.5), Inches(6.4), Inches(4.0))
    curve_tr.line.color.rgb = GOLD
    curve_tr.line.width = Pt(2)
    curve_tr.line.end_arrowhead = 2

    # --- Right Section ---

    # 1. Bottlenecks Box
    bot_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.3), Inches(1.8))
    bot_bg.fill.solid()
    bot_bg.fill.fore_color.rgb = DARK_BLUE_BG
    bot_bg.line.color.rgb = BORDER_BLUE
    bot_bg.line.width = Pt(1)

    bot_title = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(5.3), Inches(0.4))
    bot_p = bot_title.text_frame.paragraphs[0]
    bot_p.text = "四大瓶颈"
    bot_p.font.size = Pt(16)
    bot_p.font.bold = True
    bot_p.font.color.rgb = GOLD
    bot_p.alignment = PP_ALIGN.CENTER
    bot_p.font.name = FONT_NAME

    def add_list_item(left, top, icon, text):
        tb = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{icon}  {text}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(2.2), "☁", "数据字典幻觉")
    add_list_item(Inches(10.3), Inches(2.2), "🔗", "接口规则缺失")
    add_list_item(Inches(7.8), Inches(2.7), "📈", "修复成本反超")
    add_list_item(Inches(10.3), Inches(2.7), "💬", "提示词被忽略")

    # 2. Tools Box
    tool_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.6), Inches(5.3), Inches(1.8))
    tool_bg.fill.solid()
    tool_bg.fill.fore_color.rgb = DARK_BLUE_BG
    tool_bg.line.color.rgb = BORDER_BLUE
    tool_bg.line.width = Pt(1)

    tool_title = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(5.3), Inches(0.4))
    tool_p = tool_title.text_frame.paragraphs[0]
    tool_p.text = "工具选型"
    tool_p.font.size = Pt(16)
    tool_p.font.bold = True
    tool_p.font.color.rgb = GOLD
    tool_p.alignment = PP_ALIGN.CENTER
    tool_p.font.name = FONT_NAME

    # Claude Code
    cc_title = slide.shapes.add_textbox(Inches(7.8), Inches(4.1), Inches(2.5), Inches(0.4))
    cc_p = cc_title.text_frame.paragraphs[0]
    cc_p.text = "⚡ Claude Code"
    cc_p.font.size = Pt(15)
    cc_p.font.bold = True
    cc_p.font.color.rgb = WHITE
    cc_p.font.name = FONT_NAME

    add_list_item(Inches(7.8), Inches(4.5), "☑", "业务理解优秀")
    add_list_item(Inches(7.8), Inches(4.9), "☑", "多文件解析强")

    # Copilot
    cp_title = slide.shapes.add_textbox(Inches(10.3), Inches(4.1), Inches(2.5), Inches(0.4))
    cp_p = cp_title.text_frame.paragraphs[0]
    cp_p.text = "🤖 Copilot"
    cp_p.font.size = Pt(15)
    cp_p.font.bold = True
    cp_p.font.color.rgb = WHITE
    cp_p.font.name = FONT_NAME

    add_list_item(Inches(10.3), Inches(4.5), "☒", "业务理解受限")
    add_list_item(Inches(10.3), Inches(4.9), "☒", "多文件解析弱")

    # 3. Conclusion Box
    conc_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(5.7), Inches(5.3), Inches(1.4))
    conc_shadow.fill.background()
    conc_shadow.line.color.rgb = GOLD
    conc_shadow.line.width = Pt(1.5)

    conc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.6), Inches(5.3), Inches(1.4))
    conc_bg.fill.solid()
    conc_bg.fill.fore_color.rgb = RGBColor(0x1C, 0x2A, 0x35)
    conc_bg.line.color.rgb = BORDER_BLUE
    conc_bg.line.width = Pt(1)

    rb_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.8), Inches(4.9), Inches(1.1))
    rb_tf = rb_text.text_frame
    rb_tf.word_wrap = True
    rb_p = rb_tf.paragraphs[0]

    run1 = rb_p.add_run()
    run1.text = "核心建议："
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = GOLD
    run1.font.name = FONT_NAME

    run2 = rb_p.add_run()
    run2.text = "当前 AI 尚未具备处理 ABAP 中高难度任务的能力，仅推荐用于"
    run2.font.size = Pt(15)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME

    run3 = rb_p.add_run()
    run3.text = "辅助简单逻辑"
    run3.font.size = Pt(15)
    run3.font.bold = True
    run3.font.color.rgb = GOLD
    run3.font.name = FONT_NAME

    run4 = rb_p.add_run()
    run4.text = "。"
    run4.font.size = Pt(15)
    run4.font.color.rgb = WHITE
    run4.font.name = FONT_NAME
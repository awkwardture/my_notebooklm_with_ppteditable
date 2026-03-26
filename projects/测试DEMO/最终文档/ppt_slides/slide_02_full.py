from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"projects/测试DEMO/最终文档/ppt_slides/slide_02.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 定义颜色
    BG_COLOR = RGBColor(43, 43, 43)
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(170, 170, 170)
    ORANGE = RGBColor(243, 139, 0)
    PANEL_BG = RGBColor(56, 56, 56)
    PANEL_BORDER = RGBColor(80, 80, 80)
    CHART_GRAY = RGBColor(120, 120, 120)
    DARK_ORANGE_BG = RGBColor(80, 60, 40)
    DARK_GRAY_BG = RGBColor(65, 65, 65)

    # 设置背景色
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR

    # 1. 标题区域
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "中等难度：陷入“幻觉”与语法泥潭"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.6))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "案例B — 采购配额维护（函数组开发）"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"

    # 2. 左侧：图表区域
    # 图表标题
    chart_title = slide.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(2.5), Inches(0.4))
    p = chart_title.text_frame.paragraphs[0]
    p.text = "错误分布饼图"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 环形图
    chart_data = CategoryChartData()
    chart_data.categories = ['虚构字段', '其他错误']
    chart_data.add_series('Series 1', (50, 50))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(2.0), Inches(2.3), Inches(2.5), Inches(2.5), chart_data).chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = False
    
    # 设置图表颜色
    points = chart.plots[0].series[0].points
    points[0].format.fill.solid()
    points[0].format.fill.fore_color.rgb = ORANGE
    points[1].format.fill.solid()
    points[1].format.fill.fore_color.rgb = CHART_GRAY

    # 图表中心图标 (带问号的警告三角)
    center_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.05), Inches(3.3), Inches(0.4), Inches(0.35))
    center_tri.fill.background()
    center_tri.line.color.rgb = ORANGE
    center_tri.line.width = Pt(1.5)
    
    center_text = slide.shapes.add_textbox(Inches(3.05), Inches(3.35), Inches(0.4), Inches(0.35))
    p = center_text.text_frame.paragraphs[0]
    p.text = "?"
    p.font.color.rgb = ORANGE
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    center_text.text_frame.margin_top = Pt(2)

    # 图表标签 - 右侧 (橙色)
    lbl_r = slide.shapes.add_textbox(Inches(4.6), Inches(2.5), Inches(2.0), Inches(1.0))
    tf = lbl_r.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "虚构字段\n(AI Hallucinations)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "50%"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p3 = tf.add_paragraph()
    p3.text = "9处核心字段"
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY_TEXT

    # 图表标签 - 左侧 (灰色)
    lbl_l = slide.shapes.add_textbox(Inches(0.2), Inches(3.1), Inches(1.8), Inches(1.0))
    tf = lbl_l.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "其他错误\n(Other Errors)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.RIGHT
    p2 = tf.add_paragraph()
    p2.text = "21个连锁语法错误"
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.RIGHT

    # 连接线
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.8), Inches(3.4), Inches(2.1), Inches(3.4))
    line1.line.color.rgb = GRAY_TEXT
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.4), Inches(3.4), Inches(4.7), Inches(3.4))
    line2.line.color.rgb = ORANGE

    # 3. 左侧：对比表格区域
    table_top = Inches(5.2)
    table_left = Inches(0.5)
    table_width = Inches(5.5)
    table_height = Inches(2.0)

    # 外框面板
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left, table_top, table_width, table_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG
    panel.line.color.rgb = ORANGE
    panel.line.width = Pt(1)

    # 面板标题
    p_title = slide.shapes.add_textbox(table_left, table_top + Inches(0.05), table_width, Inches(0.3))
    p = p_title.text_frame.paragraphs[0]
    p.text = "AI 虚构字段 vs SAP 实际字段"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 表头
    h_y = table_top + Inches(0.4)
    h_w = Inches(2.4)
    h_h = Inches(0.25)
    
    # 左表头
    h_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(0.2), h_y, h_w, h_h)
    h_left.fill.solid()
    h_left.fill.fore_color.rgb = DARK_ORANGE_BG
    h_left.line.color.rgb = ORANGE
    p = h_left.text_frame.paragraphs[0]
    p.text = "AI 虚构"
    p.font.color.rgb = ORANGE
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 右表头
    h_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(2.9), h_y, h_w, h_h)
    h_right.fill.solid()
    h_right.fill.fore_color.rgb = PANEL_BORDER
    h_right.line.color.rgb = CHART_GRAY
    p = h_right.text_frame.paragraphs[0]
    p.text = "SAP 实际"
    p.font.color.rgb = WHITE
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 数据行
    rows_data = [
        ("Z_QUOTA_QTY", "MENG"),
        ("Z_SOURCE_VEND", "LIFNR"),
        ("Z_QUOTA_UNIT", "MEINS"),
        ("Z_VALID_TO", "DATBI"),
        ("Z_QUOTA_TYPE", "QUNUM")
    ]
    row_y_start = h_y + Inches(0.3)
    row_h = Inches(0.22)
    row_spacing = Inches(0.26)

    for i, (ai_field, sap_field) in enumerate(rows_data):
        y = row_y_start + i * row_spacing
        
        # 左侧单元格
        cell_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(0.2), y, h_w, row_h)
        cell_l.fill.solid()
        cell_l.fill.fore_color.rgb = DARK_GRAY_BG
        cell_l.line.color.rgb = ORANGE
        cell_l.line.dash_style = 4 # MSO_LINE.DASH
        p = cell_l.text_frame.paragraphs[0]
        p.text = "  " + ai_field
        p.font.color.rgb = ORANGE
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        
        # 警告图标
        icon_l = slide.shapes.add_textbox(table_left + Inches(0.2) + h_w - Inches(0.3), y - Inches(0.02), Inches(0.3), row_h)
        p_icon = icon_l.text_frame.paragraphs[0]
        p_icon.text = "⚠"
        p_icon.font.color.rgb = ORANGE
        p_icon.font.size = Pt(10)
        p_icon.alignment = PP_ALIGN.RIGHT

        # 右侧单元格
        cell_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left + Inches(2.9), y, h_w, row_h)
        cell_r.fill.solid()
        cell_r.fill.fore_color.rgb = DARK_GRAY_BG
        cell_r.line.color.rgb = CHART_GRAY
        p = cell_r.text_frame.paragraphs[0]
        p.text = "  " + sap_field
        p.font.color.rgb = WHITE
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        
        # 勾选图标
        icon_r = slide.shapes.add_textbox(table_left + Inches(2.9) + h_w - Inches(0.3), y - Inches(0.02), Inches(0.3), row_h)
        p_icon = icon_r.text_frame.paragraphs[0]
        p_icon.text = "✔"
        p_icon.font.color.rgb = GRAY_TEXT
        p_icon.font.size = Pt(10)
        p_icon.alignment = PP_ALIGN.RIGHT

    # 4. 右侧：信息面板区域
    def add_info_panel(y_pos, icon_char, title_text, desc_text, highlights=[]):
        box_w = Inches(6.2)
        box_h = Inches(1.1)
        box_x = Inches(6.5)
        
        # 面板背景
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, y_pos, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_BG
        box.line.color.rgb = PANEL_BORDER
        
        # 图标
        icon_box = slide.shapes.add_textbox(box_x + Inches(0.1), y_pos + Inches(0.15), Inches(0.8), Inches(0.8))
        p = icon_box.text_frame.paragraphs[0]
        p.text = icon_char
        p.font.size = Pt(28)
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title = slide.shapes.add_textbox(box_x + Inches(1.0), y_pos + Inches(0.1), box_w - Inches(1.1), Inches(0.3))
        p = title.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 描述文本
        desc = slide.shapes.add_textbox(box_x + Inches(1.0), y_pos + Inches(0.4), box_w - Inches(1.2), Inches(0.6))
        desc.text_frame.word_wrap = True
        p = desc.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        
        if not highlights:
            p.text = desc_text
            p.font.color.rgb = WHITE
        else:
            import re
            pattern = '(' + '|'.join(map(re.escape, highlights)) + ')'
            parts = re.split(pattern, desc_text)
            for part in parts:
                if part in highlights:
                    run = p.add_run()
                    run.text = part
                    run.font.color.rgb = ORANGE
                    run.font.bold = True
                elif part:
                    run = p.add_run()
                    run.text = part
                    run.font.color.rgb = WHITE

    # 添加四个信息面板
    add_info_panel(Inches(1.8), "🤖", "概念误导", "Copilot 完全混淆“配额”与“货源”业务模型，代码完全不可用。")
    add_info_panel(Inches(3.1), "🌩️", "严重幻觉", "Claude 虚构 9 处核心字段（占比 50%），引发 21 个连锁语法错误。", highlights=["9", "50%", "21"])
    add_info_panel(Inches(4.4), "🔒", "接口违规", "AI 无法识别 SAP FM 接口必须使用 DDIC 类型的强制规则。")
    add_info_panel(Inches(5.7), "⏱️", "效率停滞", "修正幻觉字段与重写逻辑的成本已抵消 AI 生成的便利。")

    # 5. 页码
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.4))
    p = page_num.text_frame.paragraphs[0]
    p.text = "2/4"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
